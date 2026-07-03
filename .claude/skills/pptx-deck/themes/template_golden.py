"""template_golden — 继承 tech_blue + 3 个独有 layout(process_flow / pyramid / radial)。

实现 author 用到的 3 个独有 layout:
- process_flow:N 步 vertical 流程(numbered circle + step_title + step_desc + arrows)
- pyramid:N-tier stacked structure(顶窄底宽,色梯度,可选侧栏)
- radial:central node + N spokes(圆形围绕中心,可选连线)

其他 layout(cover/toc/section_divider/single_focus/compare/cards/bullet_list/table/
pic_text/summary/closing 等 14 个)从 tech_blue re-export,共享 SSOT。

色彩 / 字体 token 来自 themes/template_golden.yaml(纯 yaml · load_theme 直接读);
此处常量跟 yaml token 对齐(apply_theme 推过来)。
"""
import math
import sys
from pathlib import Path
from typing import Any

# Fallback for direct import outside pytest
_helpers_path = str(Path(__file__).parent.parent.parent / "pptx")
if _helpers_path not in sys.path:
    sys.path.insert(0, _helpers_path)

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as _Pres
from pptx.slide import Slide
from pptx.util import Emu, Inches

import helpers as H
import layout as L

# 通过 spec_from_file_location 加载时无 parent package context,relative
# import 会失败。优先 try relative(正常 pytest / import 路径),退到
# absolute(_extract_theme_from_pptx 派生的 freestanding module)。
try:
    from . import tech_blue as _tb
except ImportError:
    _themes_dir = str(Path(__file__).parent)
    if _themes_dir not in sys.path:
        sys.path.insert(0, _themes_dir)
    import tech_blue as _tb  # type: ignore[no-redef]


# ===== Re-export tech_blue 14 个 make_*(SSOT 在 tech_blue.py) =====
make_cover = _tb.make_cover
make_toc = _tb.make_toc
make_section_divider = _tb.make_section_divider
make_single_focus = _tb.make_single_focus
make_compare = _tb.make_compare
make_compare_pk = _tb.make_compare_pk
make_matrix_2x2 = _tb.make_matrix_2x2
make_cards = _tb.make_cards
make_bullet_list = _tb.make_bullet_list
make_table = _tb.make_table
make_pic_text = _tb.make_pic_text
make_summary = _tb.make_summary
make_closing = _tb.make_closing


# ===== Token aliases(跟 template_golden.yaml token 对齐 · apply_theme 推过来) =====
FONT_HEADER = _tb.FONT_HEADER
FONT_BODY = _tb.FONT_BODY
FONT_NUM = _tb.FONT_NUM
PRIMARY = _tb.PRIMARY
PRIMARY_DEEP = _tb.PRIMARY_DEEP
PRIMARY_TINT = _tb.PRIMARY_TINT
ACCENT = _tb.ACCENT


# ===== 3 个独有 layout =====

def make_process_flow(prs: _Pres, title: str,
                       steps: list[dict[str, Any]]) -> Slide:
    """N 步 vertical 流程。

    steps: list of {title, desc}。每步渲染:左侧 numbered circle + 右侧 step_title +
    step_desc。步骤之间用淡色 connector 串连(垂直),最后一步无 connector。

    N=3-7 推荐;N>7 单页装不下,author 应拆页。
    """
    s = _tb._blank_slide(prs)
    _tb._add_title(s, title)
    region = L.content_region()
    n = len(steps)
    gap = Inches(0.18)
    row_h = Emu((region.h - gap * (n - 1)) // n)
    rows = L.stack(region, [row_h] * n, gap=gap, align="top")

    # 几何参数
    circle_d = Inches(0.7)  # numbered circle diameter
    desc_left_offset = Inches(1.0)  # step text 在 circle 右侧偏移

    for i, (row, step) in enumerate(zip(rows, steps)):
        # numbered circle(左侧)
        circle_y = row.y + Emu((row.h - circle_d) // 2)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    row.x, circle_y, circle_d, circle_d)
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        H.no_line(circle)
        # number text(circle 内居中)
        n_tb = s.shapes.add_textbox(row.x, circle_y, circle_d, circle_d)
        H.fix_textbox_margins(n_tb.text_frame)
        n_tb.text_frame.vertical_anchor = 3  # MIDDLE
        pn = n_tb.text_frame.paragraphs[0]
        pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run()
        rn.text = str(i + 1)
        H.set_font(rn, name=FONT_NUM, size=20, bold=True, color=H.WHITE)

        # connector(垂直短线,连本步到下一步;最后一步不画)
        if i < n - 1:
            conn_x = row.x + Emu(circle_d // 2) - Emu(Inches(0.015))
            conn_y = circle_y + circle_d
            conn_h = row.h - circle_d + gap
            conn = H.rect(s, conn_x, conn_y, Inches(0.03), conn_h, H.GRAY_300)
            H.no_line(conn)

        # step title + desc(circle 右侧)
        text_x = row.x + desc_left_offset
        text_w = row.w - desc_left_offset
        title_box = L.Box(text_x, row.y, text_w, Inches(0.4))
        _tb._text(s, title_box, step.get("title", ""), size=18, bold=True,
                  color=PRIMARY_DEEP, valign="middle")
        desc_box = L.Box(text_x, row.y + Inches(0.4), text_w,
                         Emu(row.h - Inches(0.4)))
        _tb._text(s, desc_box, step.get("desc", ""), size=13,
                  color=H.GRAY_700, valign="top")
    return s


def make_pyramid(prs: _Pres, title: str, tiers: list[str], *,
                 side_left: list[dict[str, str]] | None = None,
                 side_right: list[dict[str, str]] | None = None) -> Slide:
    """N-tier 金字塔。tiers 从顶到底(顶=最 high-level,底=最 detailed)。

    每 tier:横向矩形,顶 tier 最窄,底 tier 最宽,色由浅到深(顶浅底深 PRIMARY)。
    tier 内置文字白色居中。

    side_left / side_right(可选 1-2 项):各侧 stack 显 {title, body}。
    无侧栏时金字塔占中央 60% 宽,有侧栏时占中央 40% 宽。
    """
    s = _tb._blank_slide(prs)
    _tb._add_title(s, title)
    region = L.content_region()
    n = len(tiers)

    has_side = bool(side_left or side_right)
    pyramid_w_ratio = 0.42 if has_side else 0.6
    pyramid_w = Emu(int(region.w * pyramid_w_ratio))
    pyramid_x = region.x + Emu((region.w - pyramid_w) // 2)
    safe_top = Inches(0.4)
    pyramid_total_h = region.h - safe_top
    tier_gap = Inches(0.06)
    tier_h = Emu((pyramid_total_h - tier_gap * (n - 1)) // n)

    # 色梯度:tier 0(顶)= PRIMARY_TINT 浅,tier n-1(底)= PRIMARY 深
    # 线性插值 RGB
    def _lerp_color(t: float):
        from pptx.dml.color import RGBColor
        c1 = PRIMARY_TINT  # 浅
        c2 = PRIMARY       # 深
        r = int(c1[0] + (c2[0] - c1[0]) * t)
        g = int(c1[1] + (c2[1] - c1[1]) * t)
        b = int(c1[2] + (c2[2] - c1[2]) * t)
        return RGBColor(r, g, b)

    for i, tier_text in enumerate(tiers):
        # tier i 宽度:从 (1/n)*pyramid_w 线性 → pyramid_w
        # 顶 tier(i=0)最窄,底 tier(i=n-1)最宽
        width_ratio = (i + 1) / n
        this_w = Emu(int(pyramid_w * width_ratio))
        this_x = pyramid_x + Emu((pyramid_w - this_w) // 2)
        this_y = region.y + safe_top + (tier_h + tier_gap) * i
        color = _lerp_color((i + 0.5) / n)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, this_x, this_y,
                                  this_w, tier_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        H.no_line(bar)
        # tier text(白色居中)
        text_color = H.WHITE if i >= n // 2 else PRIMARY_DEEP
        text_size = 16 if H.is_handout() else 18
        t_tb = s.shapes.add_textbox(this_x, this_y, this_w, tier_h)
        H.fix_textbox_margins(t_tb.text_frame)
        t_tb.text_frame.vertical_anchor = 3
        pt = t_tb.text_frame.paragraphs[0]
        pt.alignment = PP_ALIGN.CENTER
        rt = pt.add_run()
        rt.text = tier_text
        H.set_font(rt, name=FONT_HEADER, size=text_size, bold=True,
                   color=text_color)

    # 侧栏(如有)
    if has_side:
        side_w = Emu((region.w - pyramid_w) // 2 - Inches(0.2))
        side_y = region.y + safe_top
        side_h = pyramid_total_h
        if side_left:
            _render_side_column(s, region.x, side_y, side_w, side_h,
                                side_left)
        if side_right:
            right_x = region.x + region.w - side_w
            _render_side_column(s, right_x, side_y, side_w, side_h,
                                side_right)
    return s


def _render_side_column(s: Slide, x: int, y: int, w: int, h: int,
                         items: list[dict[str, str]]) -> None:
    """侧栏:N 个 {title, body} 项,纵向 stack。"""
    n = len(items)
    gap = Inches(0.25)
    item_h = Emu((h - gap * (n - 1)) // n)
    for i, item in enumerate(items):
        item_y = y + (item_h + gap) * i
        title_box = L.Box(Emu(int(x)), Emu(int(item_y)), Emu(int(w)),
                          Inches(0.4))
        _tb._text(s, title_box, item.get("title", ""), size=14, bold=True,
                  color=PRIMARY_DEEP, valign="top")
        body_box = L.Box(Emu(int(x)), Emu(int(item_y + Inches(0.4))),
                         Emu(int(w)), Emu(int(item_h - Inches(0.4))))
        _tb._text(s, body_box, item.get("body", ""), size=12,
                  color=H.GRAY_700, valign="top")


def make_radial(prs: _Pres, title: str, center: dict[str, str],
                spokes: list[dict[str, str]]) -> Slide:
    """Central node + N spokes。spokes 围绕中心圆周排布。

    center: {title, body?}。spokes: list of {title, body}(N=3-6 推荐)。

    几何:中心圆固定位置 + N 个 spoke 圆按 angle = -90° + i*360°/N 排布。
    每个 spoke = circle + 紧邻 title + body。
    """
    s = _tb._blank_slide(prs)
    _tb._add_title(s, title)
    region = L.content_region()
    n = len(spokes)

    # 几何中心(略下移让顶 spoke 有更多上方空间放 text)
    cx = region.x + Emu(region.w // 2)
    cy = region.y + Emu(int(region.h * 0.55))
    # 中心圆 + spoke 圆尺寸
    center_d = Inches(1.4)
    spoke_d = Inches(0.9)
    # spoke 圆心到中心圆心的距离(N>=5 缩小,留 text 空间)
    radius = Inches(2.4) if n <= 4 else Inches(2.2)

    # spoke 文字尺寸(plain text 无 card box)
    spoke_text_w = Inches(1.7)
    spoke_text_h = Inches(0.7)

    # 中心圆
    cc_x = cx - Emu(center_d // 2)
    cc_y = cy - Emu(center_d // 2)
    cc = s.shapes.add_shape(MSO_SHAPE.OVAL, cc_x, cc_y, center_d, center_d)
    cc.fill.solid()
    cc.fill.fore_color.rgb = PRIMARY_DEEP
    H.no_line(cc)
    # center title(圆内居中)
    cc_tb = s.shapes.add_textbox(cc_x, cc_y, center_d, center_d)
    H.fix_textbox_margins(cc_tb.text_frame)
    cc_tb.text_frame.vertical_anchor = 3
    cc_tb.text_frame.word_wrap = True
    pc = cc_tb.text_frame.paragraphs[0]
    pc.alignment = PP_ALIGN.CENTER
    rc = pc.add_run()
    rc.text = center.get("title", "")
    H.set_font(rc, name=FONT_HEADER, size=14, bold=True, color=H.WHITE)

    # N spokes
    for i, spoke in enumerate(spokes):
        # 角度:从顶部 12 点位 (-90°) 顺时针
        angle_deg = -90 + i * (360 / n)
        angle_rad = math.radians(angle_deg)
        # spoke 圆心
        sx = cx + Emu(int(radius * math.cos(angle_rad)))
        sy = cy + Emu(int(radius * math.sin(angle_rad)))
        sc_x = sx - Emu(spoke_d // 2)
        sc_y = sy - Emu(spoke_d // 2)

        # connector line(中心圆边缘到 spoke 圆边缘)
        center_edge_x = cx + Emu(int((center_d // 2) * math.cos(angle_rad)))
        center_edge_y = cy + Emu(int((center_d // 2) * math.sin(angle_rad)))
        spoke_edge_x = sx - Emu(int((spoke_d // 2) * math.cos(angle_rad)))
        spoke_edge_y = sy - Emu(int((spoke_d // 2) * math.sin(angle_rad)))
        H.connector(s, center_edge_x, center_edge_y,
                    spoke_edge_x, spoke_edge_y, color=H.GRAY_300, weight_pt=1.5)

        # spoke 圆
        sc = s.shapes.add_shape(MSO_SHAPE.OVAL, sc_x, sc_y, spoke_d, spoke_d)
        sc.fill.solid()
        sc.fill.fore_color.rgb = PRIMARY
        H.no_line(sc)
        # spoke number(圆内居中)
        sn_tb = s.shapes.add_textbox(sc_x, sc_y, spoke_d, spoke_d)
        H.fix_textbox_margins(sn_tb.text_frame)
        sn_tb.text_frame.vertical_anchor = 3
        psn = sn_tb.text_frame.paragraphs[0]
        psn.alignment = PP_ALIGN.CENTER
        rsn = psn.add_run()
        rsn.text = str(i + 1)
        H.set_font(rsn, name=FONT_NUM, size=18, bold=True, color=H.WHITE)

        # spoke text 位置(plain text 无 card box,智能避开 spoke 圆 + region 边界)
        # 沿径向 1.55 in 偏移,确保 text 中心在 spoke 圆外(spoke_r 0.45 + gap 0.2 + text 半径 0.85)
        # top spoke(angle in [-110°, -70°])特例:text 放 spoke 下方,避开 deck title
        is_top = -110 <= angle_deg <= -70
        if is_top:
            text_cx = sx
            text_cy = sy + Emu(int(Inches(0.9)))  # spoke 圆下边缘 + 0.45 gap
        else:
            text_offset = Inches(1.55)
            text_cx = sx + Emu(int(text_offset * math.cos(angle_rad)))
            text_cy = sy + Emu(int(text_offset * math.sin(angle_rad)))

        text_x = text_cx - Emu(spoke_text_w // 2)
        text_y = text_cy - Emu(spoke_text_h // 2)
        # clamp 到 content_region 内,防出页
        text_x = max(Emu(int(region.x)),
                     min(text_x, Emu(int(region.x + region.w - spoke_text_w))))
        text_y = max(Emu(int(region.y)),
                     min(text_y, Emu(int(region.y + region.h - spoke_text_h))))

        title_box = L.Box(text_x, text_y, spoke_text_w, Inches(0.33))
        _tb._text(s, title_box, spoke.get("title", ""), size=13, bold=True,
                  color=PRIMARY_DEEP, align=PP_ALIGN.CENTER, valign="middle")
        body_box = L.Box(text_x, text_y + Inches(0.35), spoke_text_w,
                         Emu(spoke_text_h - Inches(0.35)))
        _tb._text(s, body_box, spoke.get("body", ""), size=10,
                  color=H.GRAY_700, align=PP_ALIGN.CENTER, valign="top")
    return s
