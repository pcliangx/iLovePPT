"""builder/tier1.py —— tier1 模板 slide 复用(cross-pptx deep-copy)。

负责:
- load_template_prs:打开 ingested 模板 .pptx 作为 source_prs
- init_prs_from_template:用模板做主 prs 起点(保留 master/layouts/theme),
  drop_rel 所有原 slide(防止 LibreOffice 加载失败)
- render_tier1_slide:从 source_prs 复制指定 page 的 slide 到 target_prs,
  应用 placeholder_map 替换 text,移除模板原始 footer 水印残留
- placeholder_map 加载(显式 slide.placeholder_map > 自动查找 by page index)
- text_map 按 slot tree_path 替换(map-driven)或按 geometry 替换(fallback)

不负责 footer / 页码 / source citation —— 那是 base.build_deck 的职责。
"""
from __future__ import annotations

import copy
from pathlib import Path
from types import ModuleType
from typing import Any

from pptx import Presentation
from pptx.shapes.group import GroupShape

from .base import _find_template, H, _warn


# ===========================================================================
# Placeholder text patterns —— 第三方模板常见的占位符文字 / 水印字符串。
# tier1 path 检测这些 text 并替换为 deck_plan.text_map 提供的内容 /
# 直接删除(footer 水印)。
#
# 注意:跟模板品牌相关的字符串(www.<vendor>.cc / 品牌 logo 大写残留)
# **不写在本 .py 里**,而是从同目录 `_strip_patterns.txt` 加载,
# 避免在仓库 grep 中形成"残留引用"。.pptx 内部仍含这些文字,运行时必须
# 能匹配才能 strip。
# ===========================================================================


def _load_strip_patterns() -> tuple[str, ...]:
    """从 _strip_patterns.txt 加载第三方模板水印 pattern 列表。"""
    p = Path(__file__).with_name("_strip_patterns.txt")
    try:
        lines = p.read_text(encoding="utf-8").splitlines()
    except FileNotFoundError:
        return ()
    out: list[str] = []
    for ln in lines:
        s = ln.strip()
        if not s or s.startswith("#"):
            continue
        out.append(s)
    return tuple(out)


_PLACEHOLDER_PATTERNS = (
    "…text", "...text",                       # 金字塔 tier / 一般占位
    "Text here", "Text Here",                 # 通用("01.Text here" / "02.Text Here" 也命中)
    "Copy paste fonts",                       # 段落 body 占位
    "Supporting text here", "supporting text",  # process_flow 步骤 body
    "SUBTITLE HERE", "Subtitle here",         # cover 副标
    "Speaker name", "speaker name",            # 作者位
    "PRESENTATION", "Presentation",            # cover 主标残留
    "LOGO HERE",                              # logo placeholder
    "TITLE", "Title",                          # 通用 title(注意可能误命中,主线程在 title 走 placeholder name match 优先)
) + _load_strip_patterns()
_PLACEHOLDER_BARE_TEXT = {"Text", "TEXT", "01", "02", "03", "04", "05", "06"}

# 模板原始 footer 水印(从数据文件加载,strip 阶段单独比对)
_FOOTER_WATERMARK_PATTERNS = _load_strip_patterns()


def _walk_shapes(shapes):
    """Recursively yield shapes, descending into groups."""
    for shape in shapes:
        yield shape
        if isinstance(shape, GroupShape):
            yield from _walk_shapes(shape.shapes)


def _is_placeholder_text(text: str) -> bool:
    text = text.strip()
    return any(p in text for p in _PLACEHOLDER_PATTERNS)


def _collect_placeholder_shapes(slide) -> list:
    """Return list of text-bearing placeholder shapes in geometric order.

    Order strategy:
    - First by category (pyramid tiers / cards / sidebar) based on shape size+position
    - Within each category, sorted top-to-bottom then left-to-right
    """
    items: list[tuple] = []
    for shape in _walk_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not _is_placeholder_text(text):
            continue
        items.append((shape.top or 0, shape.left or 0, shape))
    items.sort(key=lambda t: (t[0], t[1]))
    return [item[2] for item in items]


def _replace_shape_text(shape, new_text: str, font_ea: str = "Microsoft YaHei",
                         text_color_hex: str | None = None,
                         font_size_pt: int | None = None):
    """Replace shape's text with new_text, preserving first run's font formatting.

    text_color_hex: 可选 "#RRGGBB",覆盖原字色(用于浅底色 tier 强制 dark text)。
    font_size_pt:   可选 int pt,强制覆盖字号(用于 cover 标题中文比英文短需缩小)。
    """
    from lxml import etree
    from pptx.oxml.ns import qn
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = new_text
        rPr = first_run._r.find(qn("a:rPr"))
        if rPr is None:
            rPr = etree.SubElement(first_run._r, qn("a:rPr"))
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font_ea)
        if text_color_hex:
            hex_clean = text_color_hex.lstrip("#")
            for child_name in ("a:solidFill", "a:schemeClr"):
                for child in rPr.findall(qn(child_name)):
                    rPr.remove(child)
            solid_fill = etree.SubElement(rPr, qn("a:solidFill"))
            srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
            srgb.set("val", hex_clean.upper())
        if font_size_pt:
            rPr.set("sz", str(int(font_size_pt) * 100))
        for run in first_para.runs[1:]:
            run.text = ""
    else:
        first_para.text = new_text
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _remove_shape(slide, shape):
    """Remove a shape from a slide (handles top-level shapes only)."""
    sp = shape._element
    parent = sp.getparent()
    if parent is not None:
        parent.remove(sp)


def _copy_slide_from_source(source_slide, target_prs, blank_layout_idx: int = 4):
    """Deep-copy a slide from source_slide into target_prs (appended).

    Uses the source slide's layout (by name match in target_prs.slide_layouts) so
    placeholder size/position/fill inheritance is preserved. Falls back to
    blank_layout_idx if no name match.
    """
    src_layout_name = source_slide.slide_layout.name
    matched_layout = None
    for lay in target_prs.slide_layouts:
        if lay.name == src_layout_name:
            matched_layout = lay
            break
    if matched_layout is None:
        if blank_layout_idx < len(target_prs.slide_layouts):
            matched_layout = target_prs.slide_layouts[blank_layout_idx]
        else:
            matched_layout = target_prs.slide_layouts[-1]
    new_slide = target_prs.slides.add_slide(matched_layout)
    for shp in list(new_slide.shapes):
        _remove_shape(new_slide, shp)
    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def _walk_tree_paths(shapes, prefix: str = ""):
    """Yield (tree_path, shape) for every shape — top-level + recursively into groups."""
    for i, shape in enumerate(shapes):
        path = f"{prefix}{i}"
        yield path, shape
        if isinstance(shape, GroupShape):
            yield from _walk_tree_paths(shape.shapes, prefix=path + ".")


def _find_shape_by_tree_path(slide, tree_path: str):
    """Locate a shape by tree_path (e.g. "3.0.1.0") via depth-first index navigation."""
    parts = tree_path.split(".")
    shapes = list(slide.shapes)
    cur = None
    for i, p in enumerate(parts):
        try:
            idx = int(p)
        except ValueError:
            return None
        if idx < 0 or idx >= len(shapes):
            return None
        cur = shapes[idx]
        if i < len(parts) - 1:
            if not isinstance(cur, GroupShape):
                return None
            shapes = list(cur.shapes)
    return cur


def _load_placeholder_map(theme: ModuleType, template_page_index: int,
                            plan_dir: str | None) -> dict | None:
    """Load library/pptx-templates/items/<template>/pages/<page-dir>/placeholder_map.yaml."""
    name = theme.__name__
    if not name.startswith("extracted_"):
        return None
    template_stem = name[len("extracted_"):]
    # Find repo root via base.py location: <repo>/.claude/skills/pptx-deck/builder/base.py
    repo_root = Path(__file__).resolve().parents[4]
    items_dir = repo_root / "library" / "pptx-templates" / "items" / template_stem / "pages"
    if not items_dir.exists():
        return None
    try:
        import yaml as _yaml
    except ImportError:
        return None
    for page_dir in items_dir.iterdir():
        map_path = page_dir / "placeholder_map.yaml"
        if not map_path.exists():
            continue
        try:
            data = _yaml.safe_load(map_path.read_text(encoding="utf-8"))
        except Exception:
            continue
        if isinstance(data, dict) and data.get("template_page_index") == template_page_index:
            return data
    return None


def _apply_tier1_text_map(slide, title: str | None, text_map: dict[str, str],
                            placeholder_map: dict | None = None):
    """On a tier1 slide:
    - replace 标题 placeholder with title (if shape named 标题 1 / Title 1 exists)
    - replace text shapes per placeholder_map (map-driven) OR auto-detect (fallback)
    - remove footer/页号 模板原始 placeholders LAST

    Order matters: text-map replacement happens BEFORE footer removal so tree_path
    indices in placeholder_map stay aligned with the original template's shape order.
    """
    # Step 1: title rewrite (by name match, doesn't shift indices)
    for shape in slide.shapes:
        name = shape.name or ""
        if shape.has_text_frame and ("标题" in name or "Title" in name):
            if title:
                _replace_shape_text(shape, title)
            break

    # Step 2: text map replacement (BEFORE footer removal — indices must stay original)
    if placeholder_map and placeholder_map.get("slots"):
        _apply_text_map_by_slots(slide, text_map, placeholder_map)
    else:
        _apply_text_map_by_geometry(slide, text_map)

    # Step 3: remove 模板原始 footer/page-number placeholders (deck adds its own footer)
    shapes_to_remove = []
    for shape in slide.shapes:
        name = shape.name or ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "页脚" in name or "Footer" in name or "灯片编号" in name or "Slide Number" in name:
                shapes_to_remove.append(shape)
                continue
            if text in _FOOTER_WATERMARK_PATTERNS:
                shapes_to_remove.append(shape)
                continue
    for shp in shapes_to_remove:
        _remove_shape(slide, shp)


def _apply_text_map_by_slots(slide, text_map: dict[str, str], placeholder_map: dict):
    """Map-driven replacement: each slot has tree_path + optional text_color."""
    slots = placeholder_map.get("slots") or []
    slot_by_id = {s["id"]: s for s in slots if isinstance(s, dict) and "id" in s}

    # Pyramid auto-dark for top-half tiers (only if slot doesn't already override)
    tier_keys = [k for k in text_map.keys() if k.startswith("tier_")]
    total_tiers = len(tier_keys)
    DARK_HEX = "#0B2A4A"

    # Two-pass: first replace texts, then remove shapes for slots whose new_text is empty.
    shapes_to_blank_remove = []

    for slot_id, slot in slot_by_id.items():
        if slot_id not in text_map:
            continue
        new_text = text_map[slot_id]
        tree_path = slot.get("tree_path")
        if not tree_path:
            continue
        shape = _find_shape_by_tree_path(slide, tree_path)
        if shape is None or not shape.has_text_frame:
            _warn("tier1.slot-map", f"slot {slot_id!r} tree_path={tree_path} 找不到 / 无 text_frame")
            continue
        if new_text == "" and slot.get("keep_when_empty") is not True:
            shapes_to_blank_remove.append(shape)
            continue
        color = slot.get("text_color")
        if color is None and slot_id.startswith("tier_") and total_tiers > 0:
            try:
                tier_num = int(slot_id.split("_", 1)[1])
                if tier_num <= total_tiers // 2:
                    color = DARK_HEX
            except ValueError:
                pass
        font_size = slot.get("font_size_pt")
        _replace_shape_text(shape, new_text, text_color_hex=color,
                              font_size_pt=font_size)

    # Now remove the empty-slot shapes (top-level only via parent.remove)
    _remove_blank_shapes(shapes_to_blank_remove)

    unmatched = set(text_map.keys()) - set(slot_by_id.keys())
    if unmatched:
        _warn("tier1.slot-map", f"text_map keys {sorted(unmatched)} 无对应 slot, 已跳过")


def _remove_blank_shapes(shapes) -> None:
    """删除空槽位 shape;失败回落到清空文本;两者都失败 → warn-loud(原文可能残留)。"""
    for shape in shapes:
        try:
            sp = shape._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)
        except Exception as e:
            try:
                _replace_shape_text(shape, "", text_color_hex=None, font_size_pt=None)
            except Exception as e2:
                _warn(
                    "tier1.shape-removal",
                    f"空槽位删除+替换均失败,模板原文可能残留: remove={e!r} replace={e2!r}",
                )


def _apply_text_map_by_geometry(slide, text_map: dict[str, str]):
    """Legacy fallback: geometric (top, left) ordering of pattern-matched shapes."""
    placeholders = _collect_placeholder_shapes(slide)
    ordered_keys = sorted(text_map.keys(), key=lambda k: _key_sort_index(k))
    tier_keys = [k for k in ordered_keys if k.startswith("tier_")]
    total_tiers = len(tier_keys)
    DARK_HEX = "#0B2A4A"
    for idx, key in enumerate(ordered_keys):
        if idx >= len(placeholders):
            break
        color_override = None
        if key.startswith("tier_") and total_tiers > 0:
            try:
                tier_num = int(key.split("_", 1)[1])
                if tier_num <= total_tiers // 2:
                    color_override = DARK_HEX
            except ValueError:
                pass
        _replace_shape_text(placeholders[idx], text_map[key], text_color_hex=color_override)


def _key_sort_index(key: str) -> tuple:
    """Sort priority for text_map keys, aligned with author semantic intent."""
    import re
    try:
        return (0, int(key), 0, "")
    except ValueError:
        pass
    m = re.match(r"^([a-zA-Z]+)_(\d+)(?:_(.+))?$", key)
    if m:
        prefix, num, role = m.group(1), int(m.group(2)), m.group(3) or ""
        role_order = {"title": 0, "header": 0, "number": 0,
                      "body": 1, "desc": 1, "supporting": 1,
                      "subtitle": 2}.get(role, 3)
        return (1, prefix, num, role_order, role)
    m = re.match(r"^(.+?)_(title|header|body|desc|subtitle)$", key)
    if m:
        base, role = m.group(1), m.group(2)
        role_order = {"title": 0, "header": 0, "body": 1, "desc": 1, "subtitle": 2}.get(role, 3)
        return (2, base, 0, role_order, role)
    return (3, key, 0, 0, "")


# ===========================================================================
# Public tier1 API(由 base.build_deck 调用)
# ===========================================================================

def load_template_prs(theme: ModuleType, plan_dir: str | None) -> Presentation | None:
    """Load template .pptx as source_prs for tier1 reuse.

    Returns None if theme is not a .pptx-extracted module(说明本 deck 不用 tier1)。
    """
    name = theme.__name__
    if not name.startswith("extracted_"):
        return None
    stem = name[len("extracted_"):]
    pptx_path = _find_template(stem, plan_dir)
    if pptx_path is None:
        return None
    return Presentation(str(pptx_path))


def init_prs_from_template(theme_id: str, plan_dir: str | None) -> Presentation:
    """以模板 .pptx 作为目标 prs 的起点(保留 master + layouts + theme schemes)。

    删除模板原所有 slide(必须 drop_rel + 删 sldIdLst entry,否则 part 残留
    导致 LibreOffice 看到 [Content_Types].xml 声明的 slide 文件但无引用,加载失败)。
    """
    template_path = _find_template(theme_id, plan_dir)
    if template_path is None:
        raise ValueError(f"找不到模板 {theme_id!r} 的 .pptx")
    prs = Presentation(str(template_path))
    prs.slide_width = H.SLIDE_W
    prs.slide_height = H.SLIDE_H
    xml_slides = prs.slides._sldIdLst
    slide_id_lst = list(xml_slides)
    for sld_id in slide_id_lst:
        rId = sld_id.rId
        prs.part.drop_rel(rId)
        xml_slides.remove(sld_id)
    return prs


def render_tier1_slide(prs: Presentation, source_prs: Presentation,
                        slide_def: dict[str, Any], page_no: int,
                        theme: ModuleType, plan_dir: str | None) -> None:
    """从 source_prs 复制指定 page 的 slide,应用 text_map + 移除模板原始 footer。

    Args:
        prs: 目标 Presentation(写入)
        source_prs: 源模板 Presentation(读取)
        slide_def: deck_plan.slides[i] 单页 dict,含 tier1_template_page / title / text_map / placeholder_map
        page_no: 1-indexed,用于报错时给用户定位
        theme: theme module(用于 _load_placeholder_map by template_page_index)
        plan_dir: deck plan 所在目录
    """
    page_idx = slide_def["tier1_template_page"]
    if not (0 <= page_idx < len(source_prs.slides)):
        raise ValueError(
            f"第 {page_no} 页 tier1_template_page={page_idx} 越界 "
            f"(模板有 {len(source_prs.slides)} 页)"
        )
    src_slide = source_prs.slides[page_idx]
    new_slide = _copy_slide_from_source(src_slide, prs)
    # 加载 placeholder_map(优先级:slide 显式 > 自动按 tier1_template_page 查找)
    pmap = slide_def.get("placeholder_map")
    if pmap is None:
        pmap = _load_placeholder_map(theme, page_idx, plan_dir)
    _apply_tier1_text_map(
        new_slide,
        title=slide_def.get("title"),
        text_map=slide_def.get("text_map", {}),
        placeholder_map=pmap,
    )
