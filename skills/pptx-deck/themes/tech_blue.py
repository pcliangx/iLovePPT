"""iLovePPT pptx-deck — 内置科技蓝主题。

调用 [[pptx]]/helpers.py 作为底层。11 个 layout 函数对应 spec §3.4。
默认字体 Microsoft YaHei（用户决策）。
"""
import sys, warnings
from pathlib import Path
from typing import Any

# Fallback for direct import outside pytest (pytest uses pyproject.toml pythonpath)
_helpers_path = str(Path(__file__).parent.parent.parent / "pptx")
if _helpers_path not in sys.path:
    sys.path.insert(0, _helpers_path)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as _Pres
from pptx.slide import Slide
from pptx.util import Emu, Inches, Length, Pt

import helpers as H
import layout as L


# ===== 字体 / 色板 — SSOT：helpers.py =====
# tech_blue 不重新定义字体或色值,只对 helpers.py 的常量做语义别名。
# 改主题色 / 字体 → 改 helpers.py 顶部的 FONT_* / BRAND_* 常量,全 deck 联动。
FONT_HEADER = H.FONT_CN   # 默认 Microsoft YaHei
FONT_BODY   = H.FONT_CN
FONT_NUM    = H.FONT_NUM  # 默认 Helvetica Neue

PRIMARY_DEEP = H.BRAND_DARK     # #0B2A4A 深海蓝
PRIMARY      = H.BRAND_PRIMARY  # #1E6FE0 科技蓝
PRIMARY_TINT = H.BRAND_TINT     # #E6F0FC 浅蓝底
ACCENT       = H.ACCENT         # #00D1C1 青绿点睛
# 灰阶直接用 helpers.py: H.GRAY_900 / H.GRAY_700 / H.GRAY_500 / H.GRAY_300 / H.GRAY_50 / H.WHITE


def _blank_slide(prs: _Pres) -> Slide:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title(
    slide: Slide,
    text: str,
    *,
    y: Length = Inches(0.6),
    size: int = 32,
    color: RGBColor = PRIMARY_DEEP,
) -> Any:
    """页面标题 — 32pt 是行业最低实践(BCG action title 36-44pt;28pt 偏弱)。"""
    box = slide.shapes.add_textbox(Inches(0.55), y, Inches(12.2), Inches(0.9))
    tf = box.text_frame
    H.fix_textbox_margins(tf)
    r = tf.paragraphs[0].add_run()
    r.text = text
    H.set_font(r, name=FONT_HEADER, size=size, bold=True, color=color)
    return box


def _text(slide: Slide, box: "L.Box", text: str, *, size: int,
          bold: bool = False, color=None, align=PP_ALIGN.LEFT, font=None) -> None:
    """在一个 Box 内放一段文字（textbox + margin 归零 + set_font）。"""
    if color is None:
        color = H.GRAY_900
    if font is None:
        font = FONT_HEADER
    tb = slide.shapes.add_textbox(box.x, box.y, box.w, box.h)
    H.fix_textbox_margins(tb.text_frame)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    H.set_font(r, name=font, size=size, bold=bold, color=color)


def make_cover(
    prs: _Pres,
    title: str,
    subtitle: str,
    *,
    prepared_by: str = "",
    date: str = "",
    version: str = "",
    project_code: str = "",
    classification: str = "",
) -> Slide:
    """封面页:深蓝底 + 主副标 + 可选咨询稿元数据(prepared_by / date / version 等)。

    元数据布局:
    - 右上角:classification 徽标(若有,如 "CONFIDENTIAL" / "INTERNAL")
    - 左下角:prepared_by · date · version · project_code(任一非空即渲染)
    """
    s = _blank_slide(prs)
    H.rect(s, 0, 0, H.SLIDE_W, H.SLIDE_H, PRIMARY_DEEP)

    # 主副标(居中)
    region = L.full_region()
    blocks = L.stack(region, [Inches(1.4), Inches(0.8)], gap=Inches(0.3),
                     align="middle")
    _text(s, blocks[0], title, size=48, bold=True, color=H.WHITE)
    _text(s, blocks[1], subtitle, size=22, color=PRIMARY_TINT)

    # 右上 classification 徽标
    if classification:
        cls_w = Inches(2.5)
        cls_box = s.shapes.add_textbox(
            Inches(H.SLIDE_W.inches - 0.55 - 2.5), Inches(0.3), cls_w, Inches(0.35))
        H.fix_textbox_margins(cls_box.text_frame)
        cls_box.text_frame.word_wrap = False
        p = cls_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = classification.upper()
        H.set_font(r, name=FONT_HEADER, size=10, bold=True, color=PRIMARY_TINT)

    # 左下元数据 (prepared_by · date · version · project_code)
    meta_parts = [v for v in (prepared_by, date, version, project_code) if v]
    if meta_parts:
        meta_w = Inches(H.SLIDE_W.inches - 1.1)
        meta_box = s.shapes.add_textbox(
            Inches(0.55), Inches(H.SLIDE_H.inches - 0.5), meta_w, Inches(0.3))
        H.fix_textbox_margins(meta_box.text_frame)
        meta_box.text_frame.word_wrap = False
        p = meta_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = " · ".join(meta_parts)
        H.set_font(r, name=FONT_BODY, size=11, color=PRIMARY_TINT)

    return s


def make_toc(prs: _Pres, sections: list[str]) -> Slide:
    """目录页：标题"目录" + N 行章节,每行编号 + 标题。"""
    s = _blank_slide(prs)
    _add_title(s, "目录", size=42, y=Inches(0.6), color=PRIMARY_DEEP)
    rboxes = L.rows(L.content_region(), len(sections))
    for i, (rb, sec) in enumerate(zip(rboxes, sections)):
        # 序号 box（左侧固定宽度）
        num_box = L.Box(x=rb.x, y=rb.y, w=Inches(0.7), h=rb.h)
        n_tb = s.shapes.add_textbox(num_box.x, num_box.y, num_box.w, num_box.h)
        H.fix_textbox_margins(n_tb.text_frame)
        rn = n_tb.text_frame.paragraphs[0].add_run()
        rn.text = f"{i+1:02d}"
        H.set_font(rn, name=FONT_NUM, size=28, bold=True, color=PRIMARY)
        # 标题 box（剩余宽度）
        title_x = Emu(rb.x + Inches(0.9))
        title_w = Emu(rb.w - Inches(0.9))
        t_tb = s.shapes.add_textbox(title_x, rb.y, title_w, rb.h)
        H.fix_textbox_margins(t_tb.text_frame)
        rt = t_tb.text_frame.paragraphs[0].add_run()
        rt.text = sec
        H.set_font(rt, name=FONT_HEADER, size=22, color=H.GRAY_700)
    return s


def make_section_divider(prs: _Pres, num: int | str, title: str) -> Slide:
    s = _blank_slide(prs)
    region = L.full_region()
    block_h = Inches(1.6)
    band = L.stack(region, [block_h], align="middle")[0]
    # 左侧深色数字块
    num_w = Inches(1.6)
    H.rect(s, band.x, band.y, num_w, band.h, PRIMARY_DEEP)
    num_tb = s.shapes.add_textbox(band.x, band.y, num_w, band.h)
    H.fix_textbox_margins(num_tb.text_frame)
    p = num_tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rn = p.add_run()
    rn.text = str(num)
    H.set_font(rn, name=FONT_NUM, size=72, bold=True, color=H.WHITE)
    # 右侧章节标题
    title_x = Emu(band.x + num_w + Inches(0.5))
    title_w = Emu(band.w - num_w - Inches(0.5))
    title_box = L.Box(x=title_x, y=band.y, w=title_w, h=band.h)
    _text(s, title_box, title, size=44, bold=True, color=PRIMARY_DEEP)
    return s


def make_single_focus(
    prs: _Pres,
    *,
    big_text: str = "",
    big_number: str = "",
    explanation: str = "",
) -> Slide:
    s = _blank_slide(prs)
    region = L.content_region()
    blocks = L.stack(region, [Inches(1.6), Inches(0.8), Inches(0.5)],
                     gap=Inches(0.2), align="middle")
    _text(s, blocks[0], big_number, size=120, bold=True, color=PRIMARY,
          font=FONT_NUM, align=PP_ALIGN.CENTER)
    _text(s, blocks[1], big_text, size=36, bold=True, color=PRIMARY_DEEP,
          align=PP_ALIGN.CENTER)
    _text(s, blocks[2], explanation, size=18, color=H.GRAY_700,
          align=PP_ALIGN.CENTER)
    return s


def make_compare(prs: _Pres, title: str, items: list[dict[str, Any]]) -> Slide:
    """N 列对比表 — header bar 风,跟 cards 视觉拉开。

    每列 = 顶部彩色 header(主推 PRIMARY 实色+白字,其他 GRAY_300+深字)
    + 下方 body 区(主推 PRIMARY_TINT 浅填充,其他 WHITE)。
    item.recommended=True 标主推列;无 recommended 字段 = 全列等地位无高亮。
    """
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    block_h = Inches(4.6) if H.is_handout() else Inches(3.5)
    row = L.stack(region, [block_h], align="middle")[0]
    cols = L.columns(row, len(items), gap=Inches(0.15))
    header_h = Inches(0.7)
    body_size = 14 if H.is_handout() else 16
    for col, item in zip(cols, items):
        is_recommended = bool(item.get("recommended", False))
        header_fill = PRIMARY if is_recommended else H.GRAY_300
        header_color = H.WHITE if is_recommended else PRIMARY_DEEP
        body_fill = PRIMARY_TINT if is_recommended else H.WHITE
        body_color = PRIMARY_DEEP if is_recommended else H.GRAY_700

        # header 方角矩形(非 H.card,跟 cards layout 拉开)
        H.rect(s, col.x, col.y, col.w, header_h, header_fill)
        h_tb = s.shapes.add_textbox(col.x, col.y, col.w, header_h)
        H.fix_textbox_margins(h_tb.text_frame)
        h_tb.text_frame.word_wrap = True
        hp = h_tb.text_frame.paragraphs[0]
        hp.alignment = PP_ALIGN.CENTER
        hr = hp.add_run(); hr.text = item["title"]
        H.set_font(hr, name=FONT_HEADER, size=18, bold=True, color=header_color)

        # body 方角矩形 + 边框
        body_y = col.y + header_h
        body_h = block_h - header_h
        body_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, col.x, body_y,
                                         col.w, body_h)
        body_shape.fill.solid(); body_shape.fill.fore_color.rgb = body_fill
        body_shape.line.color.rgb = H.GRAY_300
        body_shape.line.width = Pt(0.75)

        # 主推 ✓ 标(右上角小角标)
        if is_recommended:
            badge_w = Inches(0.55)
            badge_x = col.x + col.w - badge_w - Inches(0.1)
            badge_y = col.y + header_h + Inches(0.1)
            badge = s.shapes.add_shape(MSO_SHAPE.OVAL, badge_x, badge_y,
                                        badge_w, badge_w)
            badge.fill.solid(); badge.fill.fore_color.rgb = ACCENT
            H.no_line(badge)
            b_tb = s.shapes.add_textbox(badge_x, badge_y, badge_w, badge_w)
            H.fix_textbox_margins(b_tb.text_frame)
            bp = b_tb.text_frame.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run(); br.text = "✓"
            H.set_font(br, name=FONT_HEADER, size=18, bold=True, color=H.WHITE)

        # body 文字
        body_box = L.Box(col.x + Inches(0.25), body_y + Inches(0.25),
                          col.w - Inches(0.5), body_h - Inches(0.5))
        _text(s, body_box, item["body"], size=body_size, color=body_color)
    return s


def make_compare_pk(
    prs: _Pres,
    title: str,
    left: dict[str, str],
    right: dict[str, str],
) -> Slide:
    """对决式两选一 — 左右两大区 + 中间巨型 VS。

    用于"二选一""新旧对决""before/after"等强对比场景。跟 make_compare(N 列对比表)
    在视觉与语义上完全不同 — pk 强调"PK 感",compare 强调"列对比"。
    """
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    vs_diameter = Inches(1.4)
    gap = Inches(0.4)
    side_w = (region.w - vs_diameter - gap * 2) // 2

    block_h = Inches(4.8) if H.is_handout() else Inches(3.8)
    block_y = region.y + (region.h - block_h) // 2

    body_size = 16 if H.is_handout() else 18

    def _side(x: int, side: dict[str, str], is_left: bool) -> None:
        # 大区底色 + 顶部 4pt accent bar
        accent_color = PRIMARY if is_left else ACCENT
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, block_y,
                                  side_w, block_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = H.GRAY_50
        bg.line.color.rgb = H.GRAY_300; bg.line.width = Pt(0.75)
        bg.adjustments[0] = 0.03
        # 顶部 accent bar
        H.rect(s, x, block_y, side_w, Inches(0.08), accent_color)
        # title + body
        title_box = L.Box(x + Inches(0.4), block_y + Inches(0.35),
                            side_w - Inches(0.8), Inches(0.9))
        _text(s, title_box, side["title"], size=28, bold=True,
              color=PRIMARY_DEEP, align=PP_ALIGN.CENTER)
        body_box = L.Box(x + Inches(0.4), block_y + Inches(1.35),
                           side_w - Inches(0.8), block_h - Inches(1.6))
        _text(s, body_box, side["body"], size=body_size,
              color=H.GRAY_700, align=PP_ALIGN.CENTER)

    _side(region.x, left, is_left=True)
    _side(region.x + side_w + gap + vs_diameter + gap, right, is_left=False)

    # 中间 VS 圆
    vs_x = region.x + side_w + gap
    vs_y = block_y + (block_h - vs_diameter) // 2
    vs_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, vs_x, vs_y,
                                     vs_diameter, vs_diameter)
    vs_circle.fill.solid(); vs_circle.fill.fore_color.rgb = PRIMARY_DEEP
    H.no_line(vs_circle)
    vs_tb = s.shapes.add_textbox(vs_x, vs_y, vs_diameter, vs_diameter)
    H.fix_textbox_margins(vs_tb.text_frame)
    vp = vs_tb.text_frame.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run(); vr.text = "VS"
    H.set_font(vr, name=FONT_HEADER, size=36, bold=True, color=H.WHITE)
    return s


def make_matrix_2x2(
    prs: _Pres,
    title: str,
    x_axis: dict[str, str],
    y_axis: dict[str, str],
    quadrants: list[dict[str, Any]],
) -> Slide:
    """BCG 2×2 经典矩阵 — 横纵轴 + 4 象限,可高亮主推象限。

    x_axis / y_axis = {low: "...", high: "..."}
    quadrants = [{pos: "tl"|"tr"|"bl"|"br", title, body, highlight}] × 4
    """
    s = _blank_slide(prs)
    _add_title(s, title)

    # 矩阵主体区域 — 留左侧/底部给轴标签
    matrix_x = Inches(2.4)
    matrix_y = Inches(1.7)
    matrix_w = Inches(10.0)
    matrix_h = Inches(4.7)
    cell_w = matrix_w // 2
    cell_h = matrix_h // 2

    body_size = 14 if H.is_handout() else 12

    positions = {
        "tl": (matrix_x, matrix_y),
        "tr": (matrix_x + cell_w, matrix_y),
        "bl": (matrix_x, matrix_y + cell_h),
        "br": (matrix_x + cell_w, matrix_y + cell_h),
    }
    for q in quadrants:
        pos = q.get("pos")
        if pos not in positions:
            raise ValueError(f"quadrant.pos 必须是 tl/tr/bl/br,得到 {pos!r}")
        qx, qy = positions[pos]
        highlight = bool(q.get("highlight", False))
        fill = PRIMARY_TINT if highlight else H.WHITE
        # 象限矩形
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, qx, qy, cell_w, cell_h)
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = PRIMARY if highlight else H.GRAY_300
        rect.line.width = Pt(1.5) if highlight else Pt(0.75)
        # title + body
        title_box = L.Box(qx + Inches(0.25), qy + Inches(0.2),
                            cell_w - Inches(0.5), Inches(0.5))
        _text(s, title_box, q["title"], size=18, bold=True,
              color=PRIMARY_DEEP if highlight else H.GRAY_900)
        body_box = L.Box(qx + Inches(0.25), qy + Inches(0.85),
                           cell_w - Inches(0.5), cell_h - Inches(1.0))
        _text(s, body_box, q.get("body", ""), size=body_size,
              color=PRIMARY_DEEP if highlight else H.GRAY_700)

    # 横轴标签(矩阵下方)
    axis_y = matrix_y + matrix_h + Inches(0.1)
    x_low = L.Box(matrix_x, axis_y, cell_w, Inches(0.35))
    x_high = L.Box(matrix_x + cell_w, axis_y, cell_w, Inches(0.35))
    _text(s, x_low, x_axis.get("low", ""), size=12, bold=True,
          color=H.GRAY_700, align=PP_ALIGN.CENTER)
    _text(s, x_high, x_axis.get("high", ""), size=12, bold=True,
          color=H.GRAY_700, align=PP_ALIGN.CENTER)

    # 纵轴标签(矩阵左侧,横排小字 — 旋转文字 PPT 实现复杂,简化处理)
    y_axis_x = Inches(0.55)
    y_axis_w = matrix_x - y_axis_x - Inches(0.15)
    y_high = L.Box(y_axis_x, matrix_y, y_axis_w, Inches(0.4))
    y_low = L.Box(y_axis_x, matrix_y + matrix_h - Inches(0.4), y_axis_w,
                    Inches(0.4))
    _text(s, y_high, "↑ " + y_axis.get("high", ""), size=12, bold=True,
          color=H.GRAY_700)
    _text(s, y_low, "↓ " + y_axis.get("low", ""), size=12, bold=True,
          color=H.GRAY_700)
    return s


def make_cards(prs: _Pres, title: str, cards: list[dict[str, str]]) -> Slide:
    """N 张并列卡片。handout mode body 字号降 + box 加高。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    card_h = Inches(4.6) if H.is_handout() else Inches(3.4)
    row = L.stack(region, [card_h], align="middle")[0]
    cols = L.columns(row, len(cards))
    body_size = 12 if H.is_handout() else 16
    body_box_h = Inches(3.6) if H.is_handout() else Inches(2.2)
    for col, card in zip(cols, cards):
        H.card(s, col.x, col.y, col.w, col.h, fill=H.WHITE,
               border=H.GRAY_300, accent=PRIMARY)
        inner = L.inset(col, Inches(0.3), Inches(0.25))
        parts = L.stack(inner, [Inches(0.6), body_box_h], gap=Inches(0.15),
                        align="top")
        _text(s, parts[0], card["title"], size=20, bold=True, color=PRIMARY_DEEP)
        _text(s, parts[1], card["body"], size=body_size, color=H.GRAY_700)
    return s


def make_bullet_list(prs: _Pres, title: str, items: list[str]) -> Slide:
    """要点列表。handout mode 字号降到 14pt + 行距加大(因长句多)"""
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    # handout mode 14pt + 行距 1.6;speaker mode 18pt + 行距 1.45
    bullet_size = 14 if H.is_handout() else 18
    line_factor = 1.6 if H.is_handout() else 1.45
    line_h = Emu(int(Pt(bullet_size) * line_factor))
    block = L.stack(region, [Emu(line_h * len(items))], align="middle")[0]
    H.bullets(s, block.x, block.y, block.w, block.h, items=items,
              size=bullet_size, accent_color=PRIMARY, body_color=H.GRAY_700)
    return s


def make_table(
    prs: _Pres,
    title: str,
    headers: list[str],
    rows: list[list[str]],
) -> Slide:
    s = _blank_slide(prs)
    _add_title(s, title)
    region = L.content_region()
    H.table_modern(s, region.x, region.y, region.w, region.h,
                   headers=headers, rows=rows,
                   header_fill=PRIMARY_DEEP, header_color=H.WHITE,
                   zebra=PRIMARY_TINT, font_size=14)
    return s


def make_pic_text(
    prs: _Pres,
    title: str,
    image_path: str,
    points: list[dict[str, str]],
) -> Slide:
    """左图右文。handout mode body 字号降 + box 加高。"""
    s = _blank_slide(prs)
    _add_title(s, title)
    left, right = L.split(L.content_region(), 0.42)
    H.embed_picture(s, image_path, left.x, left.y, box_w=left.w, box_h=left.h)
    rboxes = L.rows(right, len(points))
    body_size = 11 if H.is_handout() else 13
    body_box_h = Inches(0.9) if H.is_handout() else Inches(0.5)
    for rb, p in zip(rboxes, points):
        H.card(s, rb.x, rb.y, rb.w, rb.h, fill=H.WHITE,
               border=H.GRAY_300, accent=PRIMARY)
        inner = L.inset(rb, Inches(0.25), Inches(0.12))
        parts = L.stack(inner, [Inches(0.4), body_box_h], gap=Inches(0.05),
                        align="top")
        _text(s, parts[0], p["title"], size=16, bold=True, color=PRIMARY_DEEP)
        _text(s, parts[1], p["body"], size=body_size, color=H.GRAY_700)
    return s


def make_summary(
    prs: _Pres,
    conclusions: list[str],
    title: str = "核心结论",
) -> Slide:
    """总结 N 条结论。handout mode 字号降 + 行高加大(因结论可能 60 字)"""
    s = _blank_slide(prs)
    _add_title(s, title, size=36, color=PRIMARY_DEEP)
    rboxes = L.rows(L.content_region(), len(conclusions))
    text_size = 14 if H.is_handout() else 20
    for i, (rb, c) in enumerate(zip(rboxes, conclusions)):
        num_w = Inches(0.9)
        H.rect(s, rb.x, rb.y, num_w, rb.h, PRIMARY)
        n_tb = s.shapes.add_textbox(rb.x, rb.y, num_w, rb.h)
        H.fix_textbox_margins(n_tb.text_frame)
        pn = n_tb.text_frame.paragraphs[0]
        pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run()
        rn.text = str(i + 1)
        H.set_font(rn, name=FONT_NUM, size=32, bold=True, color=H.WHITE)
        text_x = Emu(rb.x + num_w + Inches(0.25))
        text_w = Emu(rb.w - num_w - Inches(0.25))
        t_tb = s.shapes.add_textbox(text_x, rb.y, text_w, rb.h)
        H.fix_textbox_margins(t_tb.text_frame)
        t_tb.text_frame.word_wrap = True   # handout 长结论必须换行
        rt = t_tb.text_frame.paragraphs[0].add_run()
        rt.text = c
        H.set_font(rt, name=FONT_HEADER, size=text_size, color=H.GRAY_700)
    return s


def make_closing(
    prs: _Pres,
    subtitle: str = "",
    next_steps: list[dict[str, str]] | None = None,
) -> Slide:
    """封底页。两种模式:

    - 简单模式(`next_steps=None`):大字"谢谢" + 可选 subtitle 联系方式。
    - 结构化模式(`next_steps=[{action, owner?, due?}, ...]`):
      标题"Next Steps" + 编号 action 列表 + 底部 subtitle。
      咨询稿标准 closing 应当承接 Pyramid 论证,给出可执行 action,而非
      只说"谢谢"——这是把最后一页用满的方式。
    """
    s = _blank_slide(prs)
    H.rect(s, 0, 0, H.SLIDE_W, H.SLIDE_H, PRIMARY_DEEP)

    if next_steps:
        # 结构化 closing:Next Steps 列表
        title_box = L.Box(x=Inches(0.55), y=Inches(0.7),
                          w=Emu(H.SLIDE_W - Inches(1.1)), h=Inches(0.8))
        _text(s, title_box, "Next Steps", size=36, bold=True, color=H.WHITE)

        list_y = Inches(2.0)
        line_h = Inches(0.7)
        for i, step in enumerate(next_steps):
            row_y = Emu(list_y + i * line_h)
            # 序号
            num_box = s.shapes.add_textbox(Inches(0.55), row_y,
                                            Inches(0.6), line_h)
            H.fix_textbox_margins(num_box.text_frame)
            p = num_box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = f"{i + 1}."
            H.set_font(r, name=FONT_NUM, size=22, bold=True, color=ACCENT)
            # action + owner + due
            action = step.get("action", "")
            owner = step.get("owner", "")
            due = step.get("due", "")
            tail = " · ".join([v for v in (owner, due) if v])
            text = f"{action}    [{tail}]" if tail else action
            text_box = s.shapes.add_textbox(Inches(1.2), row_y,
                                             Emu(H.SLIDE_W - Inches(1.75)), line_h)
            H.fix_textbox_margins(text_box.text_frame)
            p2 = text_box.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = text
            H.set_font(r2, name=FONT_BODY, size=18, color=H.WHITE)

        if subtitle:
            sub_box = s.shapes.add_textbox(Inches(0.55),
                                            Inches(H.SLIDE_H.inches - 0.7),
                                            Emu(H.SLIDE_W - Inches(1.1)),
                                            Inches(0.4))
            H.fix_textbox_margins(sub_box.text_frame)
            p3 = sub_box.text_frame.paragraphs[0]
            p3.alignment = PP_ALIGN.LEFT
            r3 = p3.add_run()
            r3.text = subtitle
            H.set_font(r3, name=FONT_BODY, size=14, color=PRIMARY_TINT)
    else:
        # 简单模式:谢谢 + subtitle
        region = L.full_region()
        blocks = L.stack(region, [Inches(1.5), Inches(0.6)], gap=Inches(0.3),
                         align="middle")
        _text(s, blocks[0], "谢谢", size=72, bold=True, color=H.WHITE,
              align=PP_ALIGN.CENTER)
        _text(s, blocks[1], subtitle, size=18, color=PRIMARY_TINT,
              font=FONT_BODY, align=PP_ALIGN.CENTER)

    return s
