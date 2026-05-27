import subprocess
import sys
from pathlib import Path

SCRIPT = Path(__file__).parent.parent / "extractor_self_check.py"
FIXTURES = Path(__file__).parent / "fixtures"


def run(name: str, fixture_root: Path) -> tuple[int, str]:
    """运行 self_check 脚本,返回 (exit_code, stdout+stderr)."""
    r = subprocess.run(
        [sys.executable, str(SCRIPT), name, "--items-root", str(fixture_root)],
        capture_output=True, text=True,
    )
    return r.returncode, r.stdout + r.stderr


def test_ok_fixture_passes():
    code, out = run("minimal_template_ok", FIXTURES)
    assert code == 0, f"expected exit 0, got {code}\n{out}"


def test_bad_enum_fails_with_code_1():
    code, out = run("bad_enum", FIXTURES)
    assert code == 1, f"expected exit 1, got {code}\n{out}"
    assert "ENUM_VIOLATION" in out
    assert "comparison_venn" in out


def test_bad_yaml_fails_with_code_3():
    code, out = run("bad_yaml", FIXTURES)
    assert code == 3, f"expected exit 3, got {code}\n{out}"
    assert "YAML_SYNTAX_ERROR" in out


def test_missing_fields_fails_with_code_1():
    code, out = run("missing_fields", FIXTURES)
    assert code == 1, f"expected exit 1, got {code}\n{out}"
    assert "MISSING_PAGE_FIELD" in out
    assert "keywords" in out


def test_nonexistent_template_fails_with_code_4():
    code, out = run("does_not_exist", FIXTURES)
    assert code == 4
    assert "TEMPLATE_NOT_FOUND" in out


def test_yaml_list_at_top_level_exits_3(tmp_path):
    """YAML list (not dict) at template level → exit 3, no traceback."""
    fix = tmp_path / "yaml_list"
    fix.mkdir()
    (fix / "meta.yaml.draft").write_text("- item1\n- item2\n")
    pages = fix / "pages" / "01-cover"
    pages.mkdir(parents=True)
    (pages / "meta.yaml.draft").write_text("status: draft\nlayout_type: cover\n")
    code, out = run("yaml_list", tmp_path)
    assert code == 3, f"{code}\n{out}"
    assert "YAML_NOT_A_DICT" in out


def test_null_provenance_does_not_crash(tmp_path):
    """provenance: null → caught as EMBEDDING_DIM_WRONG, not AttributeError."""
    fix = tmp_path / "null_prov"
    fix.mkdir()
    # Build minimal template with provenance: null
    (fix / "meta.yaml.draft").write_text(
        "status: draft\nid: null_prov\nname: x\ncategory: test\n"
        "content_intent: [x]\nwhen_to_use: [x]\nkeywords: [x]\n"
        "recommended_for: [x]\nvisual_signature: [x]\n"
        "provenance: null\nextraction: null\n"
    )
    pages = fix / "pages" / "01-cover"
    pages.mkdir(parents=True)
    (pages / "meta.yaml.draft").write_text(
        "status: draft\nconfidence: 0.9\nid: null_prov__01-cover\n"
        "name: x\nlayout_type: cover\ncontent_intent: [x]\nwhen_to_use: [x]\n"
        "keywords: [x]\nnative_elements: [x]\n"
    )
    code, out = run("null_prov", tmp_path)
    assert code == 1, f"{code}\n{out}"
    assert "Traceback" not in out
    assert "EMBEDDING_DIM_WRONG" in out


def test_null_layout_type_does_not_pass(tmp_path):
    """layout_type: null → LAYOUT_TYPE_INVALID, exit 1, not silent pass."""
    import shutil
    src = FIXTURES / "minimal_template_ok"
    dst = tmp_path / "null_layout"
    shutil.copytree(src, dst)
    page = dst / "pages" / "01-cover" / "meta.yaml.draft"
    page.write_text(page.read_text().replace("layout_type: cover", "layout_type: null"))
    code, out = run("null_layout", tmp_path)
    assert code == 1
    assert "LAYOUT_TYPE_INVALID" in out


def test_provenance_as_list_does_not_crash(tmp_path):
    """provenance: [a, b] (non-null non-dict) → no AttributeError, EMBEDDING_DIM_WRONG fires."""
    fix = tmp_path / "prov_list"
    fix.mkdir()
    (fix / "meta.yaml.draft").write_text(
        "status: draft\nid: prov_list\nname: x\ncategory: test\n"
        "content_intent: [x]\nwhen_to_use: [x]\nkeywords: [x]\n"
        "recommended_for: [x]\nvisual_signature: [x]\n"
        "provenance: [a, b]\nextraction: [c, d]\n"
    )
    pages = fix / "pages" / "01-cover"
    pages.mkdir(parents=True)
    (pages / "meta.yaml.draft").write_text(
        "status: draft\nconfidence: 0.9\nid: prov_list__01-cover\n"
        "name: x\nlayout_type: cover\ncontent_intent: [x]\nwhen_to_use: [x]\n"
        "keywords: [x]\nnative_elements: [x]\n"
    )
    code, out = run("prov_list", tmp_path)
    assert code == 1, f"{code}\n{out}"
    assert "Traceback" not in out
    assert "EMBEDDING_DIM_WRONG" in out


def test_pmap_tree_path_unresolvable_returns_2(tmp_path):
    """Construct a placeholder_map.draft with bad tree_path, verify exit code 2."""
    import shutil
    # Set up items/ root structure
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    # Place sample.pptx as _source/
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(
        FIXTURES / "sample.pptx",
        source_dir / "minimal_template_ok.pptx",
    )
    # Write bad placeholder_map.draft (tree_path 99 doesn't exist)
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    tree_path: '99'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 2, f"{code}\n{out}"
    assert "PMAP_TREE_PATH_UNRESOLVABLE" in out


def test_pmap_valid_tree_path_passes(tmp_path):
    """A valid tree_path should not trigger PMAP error."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(
        FIXTURES / "sample.pptx",
        source_dir / "minimal_template_ok.pptx",
    )
    # tree_path '0' is the first shape on page 0 of sample.pptx (the title textbox)
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    tree_path: '0'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 0, f"{code}\n{out}"


def test_pmap_check_reports_source_missing(tmp_path):
    """If _source/<name>.pptx doesn't exist but a placeholder_map exists,
    check #11 explicitly reports SOURCE_PPTX_MISSING (not a silent skip)."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    # No _source/ dir created
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    tree_path: '99'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    # check #11 surfaces missing source as a fail (exit 1) — not a silent pass.
    assert code == 1, f"{code}\n{out}"
    assert "SOURCE_PPTX_MISSING" in out


def test_pmap_check_skipped_when_no_pmap_and_no_source(tmp_path):
    """If no placeholder_map.yaml AND no _source/ exist, no PMAP check fires."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    # No _source/ dir, no placeholder_map.yaml anywhere.
    code, out = run("minimal_template_ok", items_root)
    assert code == 0, f"{code}\n{out}"


def test_extraction_bool_rejected_as_non_int(tmp_path):
    """declared_pages: true / rendered_pages: false → EXTRACTION_TYPE_INVALID, not silent pass."""
    import shutil
    src = FIXTURES / "minimal_template_ok"
    dst = tmp_path / "ext_bool"
    shutil.copytree(src, dst)
    tpl = dst / "meta.yaml.draft"
    tpl.write_text(
        tpl.read_text()
        .replace("declared_pages: 1", "declared_pages: true")
        .replace("rendered_pages: 1", "rendered_pages: false")
        .replace("discrepancy: 0", "discrepancy: true")
    )
    code, out = run("ext_bool", tmp_path)
    assert code == 1, f"{code}\n{out}"
    assert "EXTRACTION_TYPE_INVALID" in out


# ---- check #10: list element type ----

def test_keywords_with_int_fails_check_10(tmp_path):
    """keywords list with non-str element → LIST_ELEMENT_TYPE_FAILED, exit 1.

    Reproduces the bug we hit in this session: keywords[4] = 1 → embed_text crash.
    """
    import shutil
    src = FIXTURES / "minimal_template_ok"
    dst = tmp_path / "kw_int"
    shutil.copytree(src, dst)
    tpl = dst / "meta.yaml.draft"
    tpl.write_text(tpl.read_text().replace(
        "keywords:\n  - test",
        "keywords:\n  - test\n  - 1",
    ))
    code, out = run("kw_int", tmp_path)
    assert code == 1, f"{code}\n{out}"
    assert "LIST_ELEMENT_TYPE_FAILED" in out
    assert "keywords[1]" in out


def test_page_native_elements_with_dict_fails_check_10(tmp_path):
    """page-level native_elements containing a dict element → LIST_ELEMENT_TYPE_FAILED."""
    import shutil
    src = FIXTURES / "minimal_template_ok"
    dst = tmp_path / "ne_dict"
    shutil.copytree(src, dst)
    page = dst / "pages" / "01-cover" / "meta.yaml.draft"
    page.write_text(page.read_text().replace(
        "native_elements:\n  - 标题",
        "native_elements:\n  - 标题\n  - {nested: yes}",
    ))
    code, out = run("ne_dict", tmp_path)
    assert code == 1, f"{code}\n{out}"
    assert "LIST_ELEMENT_TYPE_FAILED" in out
    assert "native_elements[1]" in out


def test_all_string_lists_pass_check_10(tmp_path):
    """Pure-string lists in both template + page → check #10 quiet."""
    import shutil
    src = FIXTURES / "minimal_template_ok"
    dst = tmp_path / "minimal_template_ok"
    shutil.copytree(src, dst)
    code, out = run("minimal_template_ok", tmp_path)
    assert code == 0, f"{code}\n{out}"
    assert "LIST_ELEMENT_TYPE_FAILED" not in out


# ---- check #11: shape_id resolve ----

def test_shape_id_missing_treated_as_null(tmp_path):
    """If pmap slot omits shape_id → treated as null → no #11 error."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(FIXTURES / "sample.pptx", source_dir / "minimal_template_ok.pptx")
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    tree_path: '0'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 0, f"{code}\n{out}"
    assert "SHAPE_ID_RESOLVE_FAILED" not in out


def test_shape_id_bogus_value_fails_check_11(tmp_path):
    """Slot shape_id pointing to non-existent shape → SHAPE_ID_RESOLVE_FAILED, exit 1."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(FIXTURES / "sample.pptx", source_dir / "minimal_template_ok.pptx")
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    shape_id: 99999\n"
        "    tree_path: '0'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 1, f"{code}\n{out}"
    assert "SHAPE_ID_RESOLVE_FAILED" in out
    assert "99999" in out


def test_shape_id_non_int_fails_check_11(tmp_path):
    """shape_id: 'abc' (str) → SHAPE_ID_RESOLVE_FAILED (type check)."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(FIXTURES / "sample.pptx", source_dir / "minimal_template_ok.pptx")
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    shape_id: abc\n"
        "    tree_path: '0'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 1, f"{code}\n{out}"
    assert "SHAPE_ID_RESOLVE_FAILED" in out


def test_shape_id_valid_passes_check_11(tmp_path):
    """Valid shape_id matching a real shape on the slide → no error."""
    import shutil
    from pptx import Presentation
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    pptx_dst = source_dir / "minimal_template_ok.pptx"
    shutil.copy(FIXTURES / "sample.pptx", pptx_dst)
    # Compute a real shape_id from the sample slide 0
    pres = Presentation(str(pptx_dst))
    first_id = pres.slides[0].shapes[0].shape_id
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        f"template_page_index: 0\n"
        f"layout_class: cover\n"
        f"slots:\n"
        f"  - id: title\n"
        f"    shape_id: {first_id}\n"
        f"    tree_path: '0'\n"
        f"    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 0, f"{code}\n{out}"


def test_shape_id_null_allowed(tmp_path):
    """shape_id: null is explicitly allowed (data-drift escape hatch)."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(FIXTURES / "sample.pptx", source_dir / "minimal_template_ok.pptx")
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml.draft"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    shape_id: null\n"
        "    tree_path: '0'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 0, f"{code}\n{out}"


def test_shape_id_check_on_approved_yaml(tmp_path):
    """check #11 runs on placeholder_map.yaml (not just .draft) — extends scope."""
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(FIXTURES / "sample.pptx", source_dir / "minimal_template_ok.pptx")
    # Place a bogus shape_id in .yaml (NOT .draft)
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    shape_id: 88888\n"
        "    tree_path: '0'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    assert code == 1, f"{code}\n{out}"
    assert "SHAPE_ID_RESOLVE_FAILED" in out


def test_tree_path_drift_in_approved_yaml_does_not_break_check_9(tmp_path):
    """check #9 is scoped to .draft — broken tree_path in approved .yaml does NOT trip it.

    Reflects current finance_arrow reality: legacy approved templates may have
    drift; #9 fail-loud is reserved for new ingests. #11 (shape_id) catches new
    breakage on approved yaml separately.
    """
    import shutil
    items_root = tmp_path / "items"
    items_root.mkdir()
    src = FIXTURES / "minimal_template_ok"
    dst = items_root / "minimal_template_ok"
    shutil.copytree(src, dst)
    source_dir = tmp_path / "_source"
    source_dir.mkdir()
    shutil.copy(FIXTURES / "sample.pptx", source_dir / "minimal_template_ok.pptx")
    # Bogus tree_path inside placeholder_map.yaml (approved), shape_id null
    pmap_path = dst / "pages" / "01-cover" / "placeholder_map.yaml"
    pmap_path.write_text(
        "template_page_index: 0\n"
        "layout_class: cover\n"
        "slots:\n"
        "  - id: title\n"
        "    shape_id: null\n"
        "    tree_path: '99'\n"
        "    capacity_chars: 24\n"
    )
    code, out = run("minimal_template_ok", items_root)
    # No #9 alert because .yaml (not .draft), no #11 alert because shape_id null.
    assert code == 0, f"{code}\n{out}"
