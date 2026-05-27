"""Extractor Step 3.3 self-check · exit code 驱动 · 替代内嵌 bash.

Exit codes:
  0 = 全过
  1 = 字段缺失 / enum / 格式 / list element type / shape_id resolve / variant enum / slot_id enum / sha drift 问题
  2 = placeholder_map tree_path 不能 resolve(check #9)
  3 = YAML 语法错
  4 = 模板目录不存在

Checks:
  1.  REQUIRED_TEMPLATE_FIELDS  必填字段(模板级)
  2.  REQUIRED_PAGE_FIELDS      必填字段(页级)
  3.  ENUM_VIOLATION / LAYOUT_TYPE_INVALID
  4.  ID_FORMAT_INVALID / ID_DUPLICATE
  5.  CONFIDENCE_NOT_NUMERIC / CONFIDENCE_OUT_OF_RANGE
  6.  EMBEDDING_DIM_WRONG
  7.  EXTRACTION_MATH_INCONSISTENT / EXTRACTION_TYPE_INVALID
  8.  TEMPLATE_NAME_MISMATCH
  9.  PMAP_TREE_PATH_UNRESOLVABLE
  10. LIST_ELEMENT_TYPE_FAILED   list 字段每个 element 必须是 str
  11. SHAPE_ID_RESOLVE_FAILED    placeholder_map slots[].shape_id 必须能在 source pptx 里找到
  12. VARIANT_ENUM_FAILED        page meta.yaml.variant 必须 ∈ library/vocabularies/layout_variants.yaml
      VARIANT_LAYOUT_MISMATCH    vocab[variant].layout_type 必须 == meta.layout_type
  13. SLOT_ID_ENUM_FAILED        placeholder_map slots[].id 必须 ∈ library/vocabularies/slot_ids.yaml
  14. SOURCE_PPTX_SHA_DRIFT     _source/<name>.pptx sha256 必须 == meta.provenance.source_pptx_sha256
"""
from __future__ import annotations

import argparse
import hashlib
import re
import sys
from pathlib import Path

import yaml

ALLOWED_LAYOUTS = {
    "cover", "toc", "section_divider", "summary", "closing", "quote",
    "single_focus", "cards", "bullet_list", "data",
    "timeline", "pyramid", "venn", "radial", "process_flow", "quadrant", "comparison",
    "other",
}

REQUIRED_TEMPLATE_FIELDS = [
    "id", "name", "category", "content_intent", "when_to_use",
    "keywords", "recommended_for", "visual_signature",
    "status", "provenance", "extraction",
]

REQUIRED_PAGE_FIELDS = [
    "id", "name", "layout_type", "content_intent", "when_to_use",
    "keywords", "native_elements", "status", "confidence",
]

# check #10: list-of-string fields that must contain only `str` elements.
# 撞过的 bug: keywords[4] = 1 (int) → embed_text crash on `" ".join(...)`.
TEMPLATE_STR_LIST_FIELDS = [
    "content_intent", "when_to_use", "keywords", "recommended_for", "visual_signature",
]
PAGE_STR_LIST_FIELDS = [
    "content_intent", "when_to_use", "native_elements", "keywords",
]

ID_RE = re.compile(r"^[a-z0-9_-]+__\d{2}-[a-z_]+$")


def _resolve_tree_path(shapes, tree_path: str) -> bool:
    """tree_path like '3' / '3.16' / '3.16.0' · walks shape tree · True if resolvable."""
    if not tree_path:
        return False
    parts = tree_path.split(".")
    current = list(shapes)
    for i, part in enumerate(parts):
        try:
            idx = int(part)
        except ValueError:
            return False
        if idx < 0 or idx >= len(current):
            return False
        node = current[idx]
        if i == len(parts) - 1:
            return True
        try:
            current = list(node.shapes)
        except AttributeError:
            return False
    return True


def _collect_shape_ids(shapes) -> set[int]:
    """Walk every shape (including group children, picture / non-text shapes)
    and collect python-pptx sp.shape_id values. Used by check #11."""
    ids: set[int] = set()
    for shape in shapes:
        try:
            sid = shape.shape_id
            if sid is not None:
                ids.add(int(sid))
        except (AttributeError, ValueError, TypeError):
            pass
        if getattr(shape, "shape_type", None) == 6:  # group
            try:
                ids.update(_collect_shape_ids(shape.shapes))
            except (AttributeError, ValueError):
                pass
    return ids


# ─────────────────────────────────────────────────────────────────
# check #13 — slot_id enum compliance
# ─────────────────────────────────────────────────────────────────
# library/vocabularies/slot_ids.yaml 是 SSOT;这里 expand 一次 cache。

_SLOT_ID_ENUM_CACHE: set[str] | None = None
_SLOT_ID_VOCAB_PATH = Path(__file__).resolve().parents[2] / "vocabularies" / "slot_ids.yaml"


def _expand_slot_id(raw: str, n_range: list[int], sub_n_range: list[int], out: set[str]) -> None:
    """Expand 'card_N_label' → card_1_label..card_8_label.
    Supports up to two 'N' occurrences (first uses n_range, second uses sub_n_range)."""
    n_count = raw.count("N")
    if n_count == 0:
        out.add(raw)
        return
    if n_count == 1:
        for n in range(n_range[0], n_range[1] + 1):
            out.add(raw.replace("N", str(n), 1))
        return
    if n_count == 2:
        for n1 in range(n_range[0], n_range[1] + 1):
            tmp = raw.replace("N", str(n1), 1)
            for n2 in range(sub_n_range[0], sub_n_range[1] + 1):
                out.add(tmp.replace("N", str(n2), 1))
        return
    # n_count > 2: skip (no such patterns expected)


def _walk_slot_ids(node, default_range: list[int], scope_n: list[int] | None,
                   scope_sub: list[int] | None, out: set[str]) -> None:
    if isinstance(node, dict):
        n_range = node.get("n_range") if isinstance(node.get("n_range"), list) else (scope_n or default_range)
        sub_n_range = node.get("sub_n_range") if isinstance(node.get("sub_n_range"), list) else (scope_sub or default_range)
        ids = node.get("ids")
        if isinstance(ids, list):
            for raw in ids:
                if isinstance(raw, str):
                    _expand_slot_id(raw, n_range, sub_n_range, out)
        for k, v in node.items():
            if k in ("ids", "n_range", "sub_n_range", "positions"):
                continue
            if isinstance(v, list) and v and all(isinstance(x, str) for x in v):
                for raw in v:
                    _expand_slot_id(raw, n_range, sub_n_range, out)
            else:
                _walk_slot_ids(v, default_range, n_range, sub_n_range, out)
    elif isinstance(node, list):
        for x in node:
            _walk_slot_ids(x, default_range, scope_n, scope_sub, out)


def load_slot_id_enum(vocab_path: Path = _SLOT_ID_VOCAB_PATH) -> set[str]:
    """Load + expand library/vocabularies/slot_ids.yaml into a flat enum set.
    Caches across calls. Empty set on failure."""
    global _SLOT_ID_ENUM_CACHE
    if _SLOT_ID_ENUM_CACHE is not None:
        return _SLOT_ID_ENUM_CACHE
    if not vocab_path.exists():
        _SLOT_ID_ENUM_CACHE = set()
        return _SLOT_ID_ENUM_CACHE
    try:
        with open(vocab_path, encoding="utf-8") as f:
            data = yaml.safe_load(f)
    except yaml.YAMLError:
        _SLOT_ID_ENUM_CACHE = set()
        return _SLOT_ID_ENUM_CACHE
    if not isinstance(data, dict):
        _SLOT_ID_ENUM_CACHE = set()
        return _SLOT_ID_ENUM_CACHE
    schema = data.get("schema", {}) or {}
    default_range = schema.get("n_range_default", [1, 8])
    if not (isinstance(default_range, list) and len(default_range) == 2):
        default_range = [1, 8]
    enum: set[str] = set()
    skeleton = schema.get("skeleton_sentinel")
    if isinstance(skeleton, str):
        enum.add(skeleton)
    for k, v in data.items():
        if k in ("version", "schema", "meta"):
            continue
        # bare list-of-strings at top section (e.g. common: [title, subtitle, ...])
        if isinstance(v, list) and v and all(isinstance(x, str) for x in v):
            for raw in v:
                _expand_slot_id(raw, default_range, default_range, enum)
        else:
            _walk_slot_ids(v, default_range, None, None, enum)
    _SLOT_ID_ENUM_CACHE = enum
    return enum


def check_slot_id_enum_compliance(pmap: dict, file_path: Path, enum: set[str]) -> list[str]:
    """check #13 — every slots[].id must be in the SSOT enum.

    Empty enum (vocab file missing) skips check silently — extractor should never
    rely on this side-effect though; CI / dev env should have vocab present."""
    errors: list[str] = []
    if not enum:
        return errors
    if not isinstance(pmap, dict):
        return errors
    for i, slot in enumerate(pmap.get("slots", []) or []):
        if not isinstance(slot, dict):
            continue
        sid = slot.get("id")
        if sid is None:
            continue
        if not isinstance(sid, str):
            errors.append(
                f"SLOT_ID_ENUM_FAILED: {file_path}: slots[{i}].id={sid!r} "
                f"(type={type(sid).__name__}, expected str ∈ slot_ids.yaml enum)"
            )
            continue
        if sid not in enum:
            errors.append(
                f"SLOT_ID_ENUM_FAILED: {file_path}: slots[{i}].id={sid!r} "
                f"∉ slot_ids.yaml enum (扩 SSOT 或换 enum 内最接近值)"
            )
    return errors


# ─────────────────────────────────────────────────────────────────
# check #12 — variant enum compliance + variant↔layout_type 对账
# ─────────────────────────────────────────────────────────────────
# library/vocabularies/layout_variants.yaml 是 SSOT;每个 variant key
# 关联到 layout_type;这里 cache 一份 {variant: layout_type} mapping。

_VARIANT_VOCAB_CACHE: dict[str, str] | None = None
_VARIANT_VOCAB_PATH = Path(__file__).resolve().parents[2] / "vocabularies" / "layout_variants.yaml"


def load_variant_vocab(vocab_path: Path = _VARIANT_VOCAB_PATH) -> dict[str, str]:
    """Load layout_variants.yaml → {variant_name: layout_type}.
    Caches across calls. Empty dict on failure (silent skip)."""
    global _VARIANT_VOCAB_CACHE
    if _VARIANT_VOCAB_CACHE is not None:
        return _VARIANT_VOCAB_CACHE
    if not vocab_path.exists():
        _VARIANT_VOCAB_CACHE = {}
        return _VARIANT_VOCAB_CACHE
    try:
        with open(vocab_path, encoding="utf-8") as f:
            data = yaml.safe_load(f)
    except yaml.YAMLError:
        _VARIANT_VOCAB_CACHE = {}
        return _VARIANT_VOCAB_CACHE
    if not isinstance(data, dict):
        _VARIANT_VOCAB_CACHE = {}
        return _VARIANT_VOCAB_CACHE
    variants = data.get("variants", {}) or {}
    out: dict[str, str] = {}
    if isinstance(variants, dict):
        for k, v in variants.items():
            if isinstance(v, dict):
                lt = v.get("layout_type")
                if isinstance(lt, str):
                    out[k] = lt
    _VARIANT_VOCAB_CACHE = out
    return out


def check_variant_enum_compliance(meta: dict, file_path: Path, vocab: dict[str, str]) -> list[str]:
    """check #12 — page meta.yaml.variant 必须 ∈ vocab,
    且 vocab[variant].layout_type 必须 == meta.layout_type.

    variant 缺失 / null 视为允许(向后兼容 · agent 应尽量补);
    存在时必须 enum 合规 + layout_type 对账。
    Empty vocab (file missing) skips check silently."""
    errors: list[str] = []
    if not vocab or not isinstance(meta, dict):
        return errors
    variant = meta.get("variant")
    if variant is None or variant == "":
        return errors  # missing variant tolerated (back-compat)
    if not isinstance(variant, str):
        errors.append(
            f"VARIANT_ENUM_FAILED: {file_path}: variant={variant!r} "
            f"(type={type(variant).__name__}, expected str ∈ layout_variants.yaml)"
        )
        return errors
    if variant not in vocab:
        errors.append(
            f"VARIANT_ENUM_FAILED: {file_path}: variant={variant!r} "
            f"∉ layout_variants.yaml enum (扩 SSOT 或换 enum 内最接近值)"
        )
        return errors
    expected_lt = vocab[variant]
    actual_lt = meta.get("layout_type")
    if isinstance(actual_lt, str) and actual_lt and actual_lt != expected_lt:
        errors.append(
            f"VARIANT_LAYOUT_MISMATCH: {file_path}: variant={variant!r} "
            f"declares layout_type={expected_lt!r} but meta.layout_type={actual_lt!r}"
        )
    return errors


def check_list_element_types(data: dict, file_path: Path, fields: list[str]) -> list[str]:
    """check #10 — every element of each list-of-string field must be `str`.

    Returns list of error strings (empty when OK). Skips fields that are absent
    or not a list (those are caught by REQUIRED_FIELDS / other checks).
    """
    errors: list[str] = []
    if not isinstance(data, dict):
        return errors
    for fname in fields:
        v = data.get(fname)
        if not isinstance(v, list):
            continue
        for i, item in enumerate(v):
            if not isinstance(item, str):
                errors.append(
                    f"LIST_ELEMENT_TYPE_FAILED: {file_path}: {fname}[{i}]={item!r} "
                    f"(type={type(item).__name__}, expected str)"
                )
    return errors


def sha256_of(path: Path) -> str:
    """Compute sha256 of a file, streaming in 1MB chunks."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def load_yaml(path: Path) -> tuple[dict | None, str | None]:
    """返回 (data, error_msg). data is None 时 error_msg 非空."""
    try:
        with open(path, encoding="utf-8") as f:
            data = yaml.safe_load(f)
    except yaml.YAMLError as e:
        return None, str(e).split("\n")[0]
    except FileNotFoundError:
        return None, f"file not found: {path}"
    if data is None:
        return None, "file is empty or null"
    if not isinstance(data, dict):
        return None, f"YAML_NOT_A_DICT: top-level is {type(data).__name__}, expected dict"
    return data, None


def check(name: str, items_root: Path) -> int:
    tpl_dir = items_root / name
    if not tpl_dir.is_dir():
        print(f"TEMPLATE_NOT_FOUND: {tpl_dir}")
        return 4

    errors: list[str] = []
    syntax_errors: list[str] = []

    # 0. YAML 语法(模板级 + 所有页)
    tpl_meta_path = tpl_dir / "meta.yaml.draft"
    if not tpl_meta_path.exists():
        tpl_meta_path = tpl_dir / "meta.yaml"
    tpl_meta, err = load_yaml(tpl_meta_path)
    if err:
        syntax_errors.append(f"YAML_SYNTAX_ERROR: {tpl_meta_path}: {err}")

    page_paths = sorted(tpl_dir.glob("pages/*/meta.yaml.draft"))
    if not page_paths:
        page_paths = sorted(tpl_dir.glob("pages/*/meta.yaml"))
    page_metas: list[tuple[Path, dict]] = []
    for p in page_paths:
        data, err = load_yaml(p)
        if err:
            syntax_errors.append(f"YAML_SYNTAX_ERROR: {p}: {err}")
        elif data is not None:
            page_metas.append((p, data))

    if syntax_errors:
        for e in syntax_errors:
            print(e)
        return 3

    # 1. 模板级必填字段
    if tpl_meta:
        for f in REQUIRED_TEMPLATE_FIELDS:
            if f not in tpl_meta:
                errors.append(f"MISSING_TEMPLATE_FIELD: {tpl_meta_path}: {f}")

    # 2. 页级必填字段
    for p, data in page_metas:
        for f in REQUIRED_PAGE_FIELDS:
            if f not in data:
                errors.append(f"MISSING_PAGE_FIELD: {p}: {f}")

    # 3. layout_type enum (None/"" treated as invalid, not skipped)
    for p, data in page_metas:
        lt = data.get("layout_type")
        if not isinstance(lt, str) or not lt:
            errors.append(f"LAYOUT_TYPE_INVALID: {p}: layout_type={lt!r} (must be non-empty string)")
        elif lt not in ALLOWED_LAYOUTS:
            errors.append(f"ENUM_VIOLATION: {p}: layout_type={lt}")

    # 4. id 格式 + 唯一性
    ids = []
    for p, data in page_metas:
        page_id = data.get("id", "")
        if not ID_RE.match(page_id):
            errors.append(f"ID_FORMAT_INVALID: {p}: id={page_id!r} (expected <name>__<NN-slug>)")
        ids.append(page_id)
    if len(set(ids)) != len(ids):
        errors.append(f"ID_DUPLICATE: {len(ids)} pages but {len(set(ids))} unique ids")

    # 5. confidence 是数字
    for p, data in page_metas:
        c = data.get("confidence")
        if not isinstance(c, (int, float)):
            errors.append(f"CONFIDENCE_NOT_NUMERIC: {p}: confidence={c!r}")
        elif not 0.0 <= c <= 1.0:
            errors.append(f"CONFIDENCE_OUT_OF_RANGE: {p}: confidence={c}")

    # 6. provenance.embedding_dim == 1152
    if tpl_meta:
        prov = tpl_meta.get("provenance")
        if not isinstance(prov, dict):
            prov = {}
        if prov.get("embedding_dim") != 1152:
            errors.append(f"EMBEDDING_DIM_WRONG: {tpl_meta_path}: got {prov.get('embedding_dim')}, expected 1152")

    # 7. extraction 算式自洽
    if tpl_meta:
        ext = tpl_meta.get("extraction")
        if not isinstance(ext, dict):
            ext = {}
        d = ext.get("declared_pages")
        r = ext.get("rendered_pages")
        disc = ext.get("discrepancy")
        if d is not None and r is not None and disc is not None:
            if not all(isinstance(x, int) and not isinstance(x, bool) for x in (d, r, disc)):
                errors.append(f"EXTRACTION_TYPE_INVALID: {tpl_meta_path}: declared/rendered/discrepancy must be int, got ({type(d).__name__}, {type(r).__name__}, {type(disc).__name__})")
            elif d - r != disc:
                errors.append(f"EXTRACTION_MATH_INCONSISTENT: declared={d} rendered={r} discrepancy={disc}")

    # 8. template_name 跟父目录名一致
    for p, data in page_metas:
        tn = data.get("template_name")
        if tn and tn != name:
            errors.append(f"TEMPLATE_NAME_MISMATCH: {p}: template_name={tn} != {name}")

    # 14. source_pptx sha drift — _source/<name>.pptx 必须跟 meta.provenance.source_pptx_sha256 匹配
    #     模板 .pptx 源若更新而 sha 没 bump,placeholder_map.shape_id 静默失效。
    #     skip 静默情况(source pptx 不在 / 没 declared sha · 让别的 check 报)。
    if tpl_meta and isinstance(tpl_meta.get("provenance"), dict):
        prov = tpl_meta["provenance"]
        declared_sha = prov.get("source_pptx_sha256")
        source_pptx = items_root.parent / "_source" / f"{name}.pptx"
        if declared_sha and source_pptx.exists():
            try:
                actual_sha = sha256_of(source_pptx)
            except OSError as e:
                errors.append(f"SOURCE_PPTX_SHA_DRIFT: cannot read {source_pptx}: {e}")
            else:
                if actual_sha != declared_sha:
                    errors.append(
                        f"SOURCE_PPTX_SHA_DRIFT: {tpl_meta_path}: "
                        f"declared={declared_sha[:12]}... actual={actual_sha[:12]}... "
                        f"(模板 .pptx 源已变 · 必须重新 inspect placeholder_map · "
                        f"跑 inspect_placeholders.py + bump source_pptx_version)"
                    )

    # 10. list element types — 每个 list 字段每个 element 必须是 str
    # (在 #9 placeholder_map check 之前;失败也走通用 exit-1 路径)
    if tpl_meta:
        errors.extend(check_list_element_types(tpl_meta, tpl_meta_path, TEMPLATE_STR_LIST_FIELDS))
    for p, data in page_metas:
        errors.extend(check_list_element_types(data, p, PAGE_STR_LIST_FIELDS))

    # 12. variant enum compliance + variant↔layout_type 对账
    # (依赖 library/vocabularies/layout_variants.yaml;vocab 缺则 silently skip)
    variant_vocab = load_variant_vocab()
    if variant_vocab:
        for p, data in page_metas:
            errors.extend(check_variant_enum_compliance(data, p, variant_vocab))

    # 9. placeholder_map tree_path resolves (legacy scope: .draft only)
    # 11. shape_id resolves (new: applies to both .yaml AND .yaml.draft)
    # 13. slot_id enum compliance (applies to both .yaml AND .yaml.draft;
    #     不依赖 source pptx;直接对 library/vocabularies/slot_ids.yaml SSOT 校验)
    # Each check has its own error bucket and exit semantics:
    #   - tree_path errors alone → exit 2
    #   - shape_id errors → exit 1 (merged with general errors)
    #   - slot_id enum errors → exit 1 (merged with general errors)
    pmap_errors: list[str] = []
    shape_id_errors: list[str] = []
    slot_id_errors: list[str] = []
    source_pptx = items_root.parent / "_source" / f"{name}.pptx"

    # check #11 looks at every placeholder_map (approved + draft).
    pmap_paths_all = sorted(set(
        list(tpl_dir.glob("pages/*/placeholder_map.yaml"))
        + list(tpl_dir.glob("pages/*/placeholder_map.yaml.draft"))
    ))
    # check #9 keeps its legacy scope — only .draft (to avoid surfacing
    # pre-existing tree_path drift in approved templates here; that's a
    # separate cleanup sprint).
    pmap_paths_draft = sorted(tpl_dir.glob("pages/*/placeholder_map.yaml.draft"))

    # check #13 — slot_id enum compliance (independent of source pptx existence)
    slot_id_enum = load_slot_id_enum()
    if slot_id_enum:
        for p in pmap_paths_all:
            pmap, err = load_yaml(p)
            if err or not pmap:
                continue
            slot_id_errors.extend(check_slot_id_enum_compliance(pmap, p, slot_id_enum))

    if pmap_paths_all and not source_pptx.exists():
        # check #11 demands source .pptx to verify shape_id. Report it specifically.
        shape_id_errors.append(f"SOURCE_PPTX_MISSING: {source_pptx} (required for shape_id resolve check)")
    elif pmap_paths_all and source_pptx.exists():
        try:
            from pptx import Presentation as _Pres
            pres = _Pres(str(source_pptx))
            for p in pmap_paths_all:
                pmap, err = load_yaml(p)
                if err or not pmap:
                    continue
                idx = pmap.get("template_page_index")
                if idx is None or not isinstance(idx, int) or isinstance(idx, bool) or idx < 0 or idx >= len(pres.slides):
                    # PMAP_PAGE_INDEX_INVALID belongs to check #9 only if file is .draft.
                    if p in pmap_paths_draft:
                        pmap_errors.append(f"PMAP_PAGE_INDEX_INVALID: {p}: template_page_index={idx}")
                    else:
                        shape_id_errors.append(f"PMAP_PAGE_INDEX_INVALID: {p}: template_page_index={idx}")
                    continue
                slide = pres.slides[idx]
                # Pre-collect every shape_id on this slide (groups included) for check #11.
                slide_shape_ids = _collect_shape_ids(slide.shapes)
                is_draft = p in pmap_paths_draft
                for slot in pmap.get("slots", []):
                    # check #9 — tree_path resolve (draft only)
                    if is_draft:
                        tp = slot.get("tree_path", "")
                        if tp and not _resolve_tree_path(slide.shapes, tp):
                            pmap_errors.append(f"PMAP_TREE_PATH_UNRESOLVABLE: {p}: tree_path={tp!r}")
                    # check #11 — shape_id resolve (both)
                    sid = slot.get("shape_id")
                    if sid is None:
                        # null shape_id is allowed (data drift / not-yet-backfilled);
                        # tree_path remains the fallback.
                        continue
                    if not isinstance(sid, int) or isinstance(sid, bool):
                        shape_id_errors.append(
                            f"SHAPE_ID_RESOLVE_FAILED: {p}: shape_id={sid!r} (type={type(sid).__name__}, expected int)"
                        )
                        continue
                    if sid not in slide_shape_ids:
                        shape_id_errors.append(
                            f"SHAPE_ID_RESOLVE_FAILED: {p}: shape_id={sid} not found on slide {idx} "
                            f"(available count={len(slide_shape_ids)})"
                        )
        except Exception as e:
            pmap_errors.append(f"PMAP_CHECK_FAILED: {e}")

    # shape_id + slot_id errors merge into general bucket (exit 1).
    errors.extend(shape_id_errors)
    errors.extend(slot_id_errors)
    # If ONLY pmap (tree_path) errors (no other errors), exit 2 — preserves legacy semantics.
    if pmap_errors and not errors:
        for e in pmap_errors:
            print(e)
        return 2
    # If both pmap and other errors, merge into general errors → exit 1
    errors.extend(pmap_errors)

    if errors:
        for e in errors:
            print(e)
        return 1

    print(f"OK: {name} · {len(page_metas)} pages · all checks passed")
    return 0


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("name")
    ap.add_argument("--items-root", default="library/pptx-templates/items",
                    help="items/ 目录(默认 library/pptx-templates/items)")
    args = ap.parse_args()
    sys.exit(check(args.name, Path(args.items_root)))


if __name__ == "__main__":
    main()
