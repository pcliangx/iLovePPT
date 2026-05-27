#!/usr/bin/env python3
"""Detect watermarks / copyright / brand LOGO in a source .pptx (P3-22).

Scans a template .pptx and surfaces candidate watermarks:

  1. Text frames hitting copyright / URL / brand-keyword patterns.
  2. Picture shapes anchored in a slide corner (top/bot × left/right with
     margin < 100 pt) and physically small (< 150 × 150 px = ~1.04 in @ 144 dpi).

Brand keywords come from `detect_watermark_brand_config.yaml` (next to this
script). Default config ships with an empty list and an `examples` comment
block so users can opt-in to vendor names like iSlide / OfficePlus.

Output:
  Single JSON object to stdout:
    {
      "watermarks": [
        {"slide": 1, "type": "text:url", "location": "shape:Title 1",
         "content": "https://www.islide.cc"},
        {"slide": 3, "type": "picture:corner",
         "location": "shape:Picture 4 anchor=bottom-right margin=42pt size=80x80pt",
         "content": "<picture>"}
      ],
      "summary": {"count": 2, "by_type": {"text:url": 1, "picture:corner": 1}}
    }

Exit codes:
  0 = scan succeeded (regardless of watermark count)
  1 = pptx not found / parse error / config malformed
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu

# ─────────────────────────────────────────────────────────────────
# Detection thresholds (kept as module constants for easy tuning).
# ─────────────────────────────────────────────────────────────────
CORNER_MARGIN_PT = 100.0          # picture must sit within this many pt of an edge
SMALL_PICTURE_MAX_PT = 150.0      # both dimensions must be ≤ this (≈ 150 px @ 72 dpi)
URL_PATTERN = re.compile(r"(?:https?://\S+)|(?:www\.\S+)", re.IGNORECASE)
COPYRIGHT_SYMBOLS = ("©", "®", "™")
COPYRIGHT_PHRASES = (
    "all rights reserved",
    "copyright",
    "版权所有",
    "© copyright",
)

CONFIG_PATH = Path(__file__).resolve().parent / "detect_watermark_brand_config.yaml"


def _emu_to_pt(value) -> float:
    """python-pptx returns EMU (English Metric Units). 1 pt = 12700 EMU."""
    if value is None:
        return 0.0
    return float(Emu(value).pt)


def load_brand_keywords(config_path: Path = CONFIG_PATH) -> list[str]:
    """Read brand-keyword list from yaml. Empty list on missing file."""
    if not config_path.exists():
        return []
    try:
        with open(config_path, encoding="utf-8") as f:
            data = yaml.safe_load(f)
    except yaml.YAMLError as e:
        print(f"BRAND_CONFIG_INVALID: {config_path}: {e}", file=sys.stderr)
        sys.exit(1)
    if not isinstance(data, dict):
        return []
    brands = data.get("brands") or []
    if not isinstance(brands, list):
        return []
    return [str(x) for x in brands if x]


def _shape_label(shape) -> str:
    """Best-effort human label: shape.name + shape_id fallback."""
    try:
        nm = shape.name
    except AttributeError:
        nm = None
    if nm:
        return nm
    try:
        return f"shape_id={shape.shape_id}"
    except (AttributeError, ValueError, TypeError):
        return "<unknown>"


def _scan_text_frame(text: str, brand_keywords: list[str]) -> list[tuple[str, str]]:
    """Return list of (type, matched_substring) hits for one text frame."""
    if not text or not text.strip():
        return []
    hits: list[tuple[str, str]] = []
    lower = text.lower()

    # URL pattern
    for m in URL_PATTERN.finditer(text):
        hits.append(("text:url", m.group(0)))

    # Copyright symbols
    for sym in COPYRIGHT_SYMBOLS:
        if sym in text:
            hits.append(("text:copyright_symbol", sym))
            break  # one symbol hit is enough per frame

    # Copyright phrases
    for phrase in COPYRIGHT_PHRASES:
        if phrase in lower:
            hits.append(("text:copyright_phrase", phrase))
            break

    # Brand keywords (case-insensitive substring)
    for kw in brand_keywords:
        if kw.lower() in lower:
            hits.append(("text:brand_keyword", kw))

    return hits


def _corner_anchor(shape, slide_w_pt: float, slide_h_pt: float) -> str | None:
    """Return 'top-left'/'top-right'/'bottom-left'/'bottom-right' if shape sits
    within CORNER_MARGIN_PT of one corner pair; else None."""
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    left_pt = _emu_to_pt(shape.left)
    top_pt = _emu_to_pt(shape.top)
    right_gap_pt = slide_w_pt - left_pt - _emu_to_pt(shape.width)
    bottom_gap_pt = slide_h_pt - top_pt - _emu_to_pt(shape.height)

    near_left = left_pt < CORNER_MARGIN_PT
    near_right = right_gap_pt < CORNER_MARGIN_PT
    near_top = top_pt < CORNER_MARGIN_PT
    near_bottom = bottom_gap_pt < CORNER_MARGIN_PT

    if near_top and near_left:
        return "top-left"
    if near_top and near_right:
        return "top-right"
    if near_bottom and near_left:
        return "bottom-left"
    if near_bottom and near_right:
        return "bottom-right"
    return None


def _walk_shapes(shapes, callback):
    """Depth-first walk including group children. callback(shape) per leaf."""
    for shape in shapes:
        if shape.shape_type == 6:  # group
            try:
                _walk_shapes(shape.shapes, callback)
                continue
            except (AttributeError, ValueError):
                pass
        callback(shape)


def scan_slide(slide, slide_idx: int, slide_w_pt: float, slide_h_pt: float,
               brand_keywords: list[str]) -> list[dict]:
    """Return list of watermark records for one slide (1-indexed slide_idx)."""
    findings: list[dict] = []

    def visit(shape):
        # Text scan
        try:
            has_tf = shape.has_text_frame
        except AttributeError:
            has_tf = False
        if has_tf:
            try:
                text = shape.text_frame.text or ""
            except AttributeError:
                text = ""
            for wtype, content in _scan_text_frame(text, brand_keywords):
                findings.append({
                    "slide": slide_idx,
                    "type": wtype,
                    "location": f"shape:{_shape_label(shape)}",
                    "content": content,
                })

        # Picture scan
        # shape_type 13 = PICTURE per python-pptx MSO_SHAPE_TYPE enum
        try:
            stype = shape.shape_type
        except AttributeError:
            stype = None
        if stype == 13:
            corner = _corner_anchor(shape, slide_w_pt, slide_h_pt)
            if corner is None:
                return
            w_pt = _emu_to_pt(shape.width)
            h_pt = _emu_to_pt(shape.height)
            if w_pt > SMALL_PICTURE_MAX_PT or h_pt > SMALL_PICTURE_MAX_PT:
                return
            # corner-anchored small picture → candidate watermark
            left_pt = _emu_to_pt(shape.left)
            top_pt = _emu_to_pt(shape.top)
            right_gap_pt = slide_w_pt - left_pt - w_pt
            bottom_gap_pt = slide_h_pt - top_pt - h_pt
            margin_pt = min(left_pt, top_pt, right_gap_pt, bottom_gap_pt)
            findings.append({
                "slide": slide_idx,
                "type": "picture:corner",
                "location": (
                    f"shape:{_shape_label(shape)} anchor={corner} "
                    f"margin={margin_pt:.0f}pt size={w_pt:.0f}x{h_pt:.0f}pt"
                ),
                "content": "<picture>",
            })

    _walk_shapes(slide.shapes, visit)
    return findings


def detect(pptx_path: Path, brand_keywords: list[str]) -> dict:
    if not pptx_path.exists():
        print(f"PPTX_NOT_FOUND: {pptx_path}", file=sys.stderr)
        sys.exit(1)
    try:
        pres = Presentation(str(pptx_path))
    except Exception as e:
        print(f"PPTX_PARSE_ERROR: {e}", file=sys.stderr)
        sys.exit(1)

    slide_w_pt = _emu_to_pt(pres.slide_width)
    slide_h_pt = _emu_to_pt(pres.slide_height)

    all_findings: list[dict] = []
    for i, slide in enumerate(pres.slides, start=1):
        all_findings.extend(scan_slide(slide, i, slide_w_pt, slide_h_pt, brand_keywords))

    by_type: dict[str, int] = {}
    for f in all_findings:
        by_type[f["type"]] = by_type.get(f["type"], 0) + 1

    return {
        "watermarks": all_findings,
        "summary": {
            "count": len(all_findings),
            "by_type": by_type,
            "pptx_path": str(pptx_path),
            "brand_keywords_used": brand_keywords,
        },
    }


def main():
    ap = argparse.ArgumentParser(description="Detect watermarks in a .pptx file")
    ap.add_argument("pptx_path", type=Path)
    ap.add_argument("--config", type=Path, default=CONFIG_PATH,
                    help="brand-keyword yaml (default: detect_watermark_brand_config.yaml next to script)")
    args = ap.parse_args()
    brand_keywords = load_brand_keywords(args.config)
    result = detect(args.pptx_path, brand_keywords)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
