"""scripts/check_source_fidelity.py claim 级保真校验测试。"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest
import yaml

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "scripts"))

import check_source_fidelity as csf  # noqa: E402

CONTENT_MD = """# Deck 标题

## 1. 营收概览
<!-- layout: single_focus -->
Q1 营收 120,000 万元,同比 +18%。

## 3. 渠道扩张
<!-- layout: cards -->
新增门店 4,500 家。

## 5. 盈利能力
<!-- layout: data -->
毛利率 38.5%,环比改善。
"""


@pytest.fixture()
def content_md(tmp_path) -> Path:
    p = tmp_path / "deck_v1_content.md"
    p.write_text(CONTENT_MD, encoding="utf-8")
    return p


def _claims_file(tmp_path, claims: list[dict]) -> Path:
    p = tmp_path / "deck_v1_claims.yaml"
    p.write_text(yaml.safe_dump({"claims": claims}, allow_unicode=True), encoding="utf-8")
    return p


def test_pass_with_placement(content_md, tmp_path):
    claims = _claims_file(tmp_path, [
        {"id": "rev", "desc": "Q1 营收", "patterns": ["120000万"], "expect_pages": [1]},
        {"id": "margin", "patterns": ["38.5%"], "expect_pages": [5]},
    ])
    report = csf.check_claims(csf.load_claims(claims), csf.extract_pages(content_md))
    assert report["summary"]["verdict"] == "pass"
    assert all(r["status"] == "pass" for r in report["claims"])


def test_fullwidth_normalization(content_md, tmp_path):
    # 全角数字 pattern 应命中半角原文(NFKC 归一)
    claims = _claims_file(tmp_path, [
        {"id": "rev_fw", "patterns": ["１２００００万"], "expect_pages": [1]},
    ])
    report = csf.check_claims(csf.load_claims(claims), csf.extract_pages(content_md))
    assert report["claims"][0]["status"] == "pass"


def test_missing_required_fails(content_md, tmp_path, capsys):
    claims = _claims_file(tmp_path, [{"id": "ghost", "patterns": ["8.8亿"]}])
    assert csf.main([str(content_md), "--claims", str(claims)]) == 1
    out = capsys.readouterr().out
    assert '"missing"' in out


def test_misplaced(content_md, tmp_path):
    claims = _claims_file(tmp_path, [
        {"id": "margin_wrong_page", "patterns": ["38.5%"], "expect_pages": [2]},
    ])
    report = csf.check_claims(csf.load_claims(claims), csf.extract_pages(content_md))
    assert report["claims"][0]["status"] == "misplaced"
    assert report["summary"]["verdict"] == "fail"


def test_optional_missing_still_passes(content_md, tmp_path):
    claims = _claims_file(tmp_path, [
        {"id": "rev", "patterns": ["120000万"]},
        {"id": "nice_to_have", "patterns": ["不存在的数"], "required": False},
    ])
    report = csf.check_claims(csf.load_claims(claims), csf.extract_pages(content_md))
    assert report["summary"]["verdict"] == "pass"
    assert report["summary"]["missing"] == 1


def test_regex_claim(content_md, tmp_path):
    claims = _claims_file(tmp_path, [
        {"id": "margin_re", "regex": r"毛利率.{0,5}38\.5", "expect_pages": [5]},
    ])
    report = csf.check_claims(csf.load_claims(claims), csf.extract_pages(content_md))
    assert report["claims"][0]["status"] == "pass"


def test_pptx_input(tmp_path):
    pptx_mod = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx_mod.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "Q1 营收 1.2 亿元"
    deck = tmp_path / "deck.pptx"
    prs.save(deck)

    claims = _claims_file(tmp_path, [
        {"id": "rev", "patterns": ["1.2亿"], "expect_pages": [1]},
    ])
    report = csf.check_claims(csf.load_claims(claims), csf.extract_pages(deck))
    assert report["claims"][0]["status"] == "pass"
    assert report["claims"][0]["hit_pages"] == [1]


def test_schema_failloud(tmp_path, content_md):
    dup = _claims_file(tmp_path, [
        {"id": "a", "patterns": ["x"]},
        {"id": "a", "patterns": ["y"]},
    ])
    with pytest.raises(ValueError, match="重复"):
        csf.load_claims(dup)

    empty = _claims_file(tmp_path, [{"id": "b"}])
    with pytest.raises(ValueError, match="至少给一个"):
        csf.load_claims(empty)

    plain_md = tmp_path / "plain.md"
    plain_md.write_text("没有页标题的文档", encoding="utf-8")
    with pytest.raises(ValueError, match="content.md"):
        csf.extract_pages(plain_md)
