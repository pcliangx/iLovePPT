"""scripts/fix_ea_fonts.py 产物端 EA 修复测试。

链路验证:latin-only CJK deck(audit ERROR)→ fix → audit 0 ERROR。
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "scripts"))

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

import audit_pptx  # noqa: E402
import fix_ea_fonts  # noqa: E402


def _deck(tmp_path, runs: list[tuple[str, str | None]]) -> Path:
    """每个 (text, latin_font) 一页;latin_font=None 则不写字体。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for text, latin in runs:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = text
        if latin:
            run.font.name = latin  # python-pptx 只写 <a:latin> — 经典 bug 路径
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def _font_errors(path: Path) -> int:
    return audit_pptx.audit(path, ["fonts"])["fonts"]["summary"]["errors"]


def test_latin_only_cjk_fixed_to_zero_errors(tmp_path):
    deck = _deck(tmp_path, [("季度营收增长", "Helvetica"), ("毛利率提升", "Arial")])
    assert _font_errors(deck) == 2
    out = tmp_path / "fixed.pptx"
    stats = fix_ea_fonts.fix_pptx(deck, out)
    assert sum(stats.values()) == 2
    assert _font_errors(out) == 0


def test_western_latin_gets_default_ea(tmp_path):
    deck = _deck(tmp_path, [("中文内容", "Helvetica")])
    out = tmp_path / "fixed.pptx"
    fix_ea_fonts.fix_pptx(deck, out)
    findings = audit_pptx.audit(out, ["fonts"])["fonts"]
    assert findings["summary"]["ok"] == 1
    # ea 应是默认 Microsoft YaHei(latin 是西文字体,不复用)
    import zipfile
    from lxml import etree
    with zipfile.ZipFile(out) as z:
        xml = z.read("ppt/slides/slide1.xml")
    root = etree.fromstring(xml)
    ea = root.find(f".//{{{fix_ea_fonts.A_NS}}}ea")
    assert ea.get("typeface") == "Microsoft YaHei"
    cs = root.find(f".//{{{fix_ea_fonts.A_NS}}}cs")
    assert cs.get("typeface") == "Microsoft YaHei"


def test_cjk_named_latin_reused_as_ea(tmp_path):
    deck = _deck(tmp_path, [("中文内容", "Microsoft YaHei")])
    out = tmp_path / "fixed.pptx"
    fix_ea_fonts.fix_pptx(deck, out)
    import zipfile
    from lxml import etree
    with zipfile.ZipFile(out) as z:
        root = etree.fromstring(z.read("ppt/slides/slide1.xml"))
    ea = root.find(f".//{{{fix_ea_fonts.A_NS}}}ea")
    assert ea.get("typeface") == "Microsoft YaHei"


def test_pure_western_run_untouched(tmp_path):
    deck = _deck(tmp_path, [("Revenue Growth", "Helvetica")])
    out = tmp_path / "fixed.pptx"
    stats = fix_ea_fonts.fix_pptx(deck, out)
    assert stats == {}


def test_undeclared_run_untouched(tmp_path):
    """无 rPr 字体声明的 run 走继承链(audit INFO/WARNING),不属于 latin-only bug。"""
    deck = _deck(tmp_path, [("中文内容", None)])
    out = tmp_path / "fixed.pptx"
    stats = fix_ea_fonts.fix_pptx(deck, out)
    assert stats == {}


def test_inplace_creates_backup(tmp_path, capsys):
    deck = _deck(tmp_path, [("季度营收", "Arial")])
    rc = fix_ea_fonts.main([str(deck)])
    capsys.readouterr()
    assert rc == 0
    assert deck.with_suffix(".pre_ea_fix.pptx").exists()
    assert _font_errors(deck) == 0


def test_idempotent(tmp_path):
    deck = _deck(tmp_path, [("季度营收", "Arial")])
    out1 = tmp_path / "f1.pptx"
    out2 = tmp_path / "f2.pptx"
    fix_ea_fonts.fix_pptx(deck, out1)
    stats2 = fix_ea_fonts.fix_pptx(out1, out2)
    assert stats2 == {}  # 已修过的 run 不再动
