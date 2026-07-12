"""scripts/audit_pptx.py section_geometry 机械几何审计测试。

夹具构造 4 类几何状态的真实 .pptx:
- off_canvas:含文字 textbox 越出画布右缘 → WARNING
- text_overlap:两个文字 textbox bbox 相交 > 35% → WARNING
- title_alignment:≥3 页左锚标题,1 页字号偏离众数 → WARNING
- 装饰豁免:无文字 oval 出血 / ≥100pt 大数字 不报
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "scripts"))

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

import audit_pptx  # noqa: E402


def _textbox(slide, x, y, w, h, text, size_pt=18):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    return tb


@pytest.fixture()
def geometry_deck(tmp_path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # slide 1:正常标题 32pt + 无文字装饰 oval 出血(不应报)
    s1 = prs.slides.add_slide(blank)
    _textbox(s1, 0.55, 0.6, 8, 0.9, "第一章标题", 32)
    s1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(-2), Inches(5), Inches(5))

    # slide 2:正常标题 + 含文字 textbox 越出右缘(应报 off_canvas)
    s2 = prs.slides.add_slide(blank)
    _textbox(s2, 0.55, 0.6, 8, 0.9, "第二章标题", 32)
    _textbox(s2, 12.0, 3.0, 3.0, 1.0, "越界文字")

    # slide 3:标题字号 40pt 偏离众数(应报 title_alignment)+ 两个重叠 textbox
    s3 = prs.slides.add_slide(blank)
    _textbox(s3, 0.55, 0.6, 8, 0.9, "第三章标题", 40)
    _textbox(s3, 2.0, 3.0, 4.0, 1.5, "文字甲")
    _textbox(s3, 2.5, 3.2, 4.0, 1.5, "文字乙")  # 与甲相交远超 35%

    # slide 4:400pt 大数字 + 其 bbox 内的小字(装饰豁免,不应报 overlap)
    s4 = prs.slides.add_slide(blank)
    _textbox(s4, 0.55, 0.6, 8, 0.9, "第四章标题", 32)
    _textbox(s4, 2.0, 1.5, 9.0, 5.0, "87%", 400)
    _textbox(s4, 4.0, 5.0, 5.0, 0.8, "覆盖率提升到")

    out = tmp_path / "geometry_deck.pptx"
    prs.save(str(out))
    return out


@pytest.fixture()
def report(geometry_deck):
    return audit_pptx.audit(geometry_deck, ["geometry"])["geometry"]


def test_off_canvas_text_flagged(report):
    hits = [f for f in report["findings"] if f["check"] == "off_canvas"]
    assert len(hits) == 1
    assert hits[0]["slide"] == 2
    assert "越界文字" in hits[0]["text"]


def test_decor_oval_bleed_not_flagged(report):
    assert not any(f["check"] == "off_canvas" and f["slide"] == 1
                   for f in report["findings"])


def test_text_overlap_flagged(report):
    hits = [f for f in report["findings"] if f["check"] == "text_overlap"]
    assert len(hits) == 1
    assert hits[0]["slide"] == 3


def test_big_number_not_flagged_as_overlap(report):
    assert not any(f["check"] == "text_overlap" and f["slide"] == 4
                   for f in report["findings"])


def test_title_size_drift_flagged(report):
    hits = [f for f in report["findings"] if f["check"] == "title_alignment"]
    assert len(hits) == 1
    assert hits[0]["slide"] == 3
    assert "40pt" in hits[0]["note"]


def test_geometry_is_advisory_no_exit_effect(geometry_deck, capsys):
    """geometry 有 WARNING 也不影响 exit code(fonts 仍是唯一 gate)。"""
    rc = audit_pptx.main([str(geometry_deck), "--sections", "geometry"])
    capsys.readouterr()
    assert rc == 0


def test_summary_counts(report):
    assert report["summary"]["warnings"] == 3
    assert report["summary"]["by_check"] == {
        "off_canvas": 1, "text_overlap": 1, "title_alignment": 1}


def test_group_children_use_group_bbox_not_child_coords(tmp_path):
    """组合形状(grpSp)子形状的 a:off 是组内子坐标系 —— 几何检查必须用组
    bbox(slide 坐标),不得把子坐标当绝对坐标产生假 off_canvas。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _textbox(s, 0.55, 0.6, 8, 0.9, "组合形状页标题", 32)
    grp = s.shapes.add_group_shape()
    tb = grp.shapes.add_textbox(Inches(12.0), Inches(3.0), Inches(3.0), Inches(1.0))
    tb.text_frame.paragraphs[0].add_run().text = "组内文字"
    # 组整体移回画布内;子形状 a:off 保持组内坐标(12,3)不变
    grp.left, grp.top = Inches(1), Inches(3)
    out = tmp_path / "grouped.pptx"
    prs.save(str(out))

    boxes = audit_pptx.PptxAudit(out)._shape_boxes()[1]
    grp_boxes = [b for b in boxes if b["kind"] == "grpSp"]
    assert len(grp_boxes) == 1
    assert "组内文字" in grp_boxes[0]["text"]
    assert abs(grp_boxes[0]["x"] - 1.0) < 0.05
    # 子形状不得以组内原始坐标单独入列(修前:x=12 被当 slide 绝对坐标)
    assert not any(b["kind"] == "sp" and "组内文字" in b["text"] for b in boxes)

    report = audit_pptx.audit(out, ["geometry"])["geometry"]
    assert not any(f["check"] == "off_canvas" for f in report["findings"]), \
        "组内子形状坐标不应触发假 off_canvas"
