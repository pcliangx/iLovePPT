"""scripts/audit_pptx.py 读侧审计测试。

夹具构造 3 种 EA 字体状态的真实 .pptx:
- slide 1: run.font.name 只写 <a:latin>(经典 bug)→ ERROR
- slide 2: helpers.set_font 写全 <a:latin>+<a:ea>+<a:cs> → OK(不进 findings)
- slide 3: run 无任何字体声明 + python-pptx 默认 theme ea 为空 → WARNING
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "scripts"))

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

import audit_pptx  # noqa: E402

import helpers as H  # noqa: E402  (pyproject pythonpath: .claude/skills/pptx)


def _add_cjk_run(slide, text: str):
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    return run


@pytest.fixture()
def mixed_deck(tmp_path) -> Path:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    r1 = _add_cjk_run(s1, "季度营收增长")
    r1.font.name = "Arial"  # python-pptx 默认路径:只写 <a:latin>

    s2 = prs.slides.add_slide(blank)
    r2 = _add_cjk_run(s2, "正确的中文字体")
    H.set_font(r2, size=18)

    s3 = prs.slides.add_slide(blank)
    r3 = _add_cjk_run(s3, "无声明继承链")
    r3.hyperlink.address = "https://example.com/data"

    prs.core_properties.title = "审计夹具"
    path = tmp_path / "mixed.pptx"
    prs.save(path)
    return path


@pytest.fixture()
def clean_deck(tmp_path) -> Path:
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    r = _add_cjk_run(s1, "全部走正确路径")
    H.set_font(r, size=18)
    path = tmp_path / "clean.pptx"
    prs.save(path)
    return path


def test_fonts_classic_bug_is_error(mixed_deck):
    report = audit_pptx.audit(mixed_deck, ["fonts"])
    summary = report["fonts"]["summary"]
    assert summary["slides"] == 3
    assert summary["errors"] == 1

    errors = [f for f in report["fonts"]["findings"] if f["severity"] == "ERROR"]
    assert len(errors) == 1
    err = errors[0]
    assert err["slide"] == 1
    assert err["latin"] == "Arial"
    assert err["ea"] == ""


def test_fonts_set_font_run_is_ok(mixed_deck):
    report = audit_pptx.audit(mixed_deck, ["fonts"])
    # slide 2 的 run 显式带 ea → OK,不进 findings
    assert not [f for f in report["fonts"]["findings"] if f["slide"] == 2]
    assert report["fonts"]["summary"]["ok"] >= 1


def test_fonts_undeclared_run_is_warning_not_error(mixed_deck):
    report = audit_pptx.audit(mixed_deck, ["fonts"])
    s3 = [f for f in report["fonts"]["findings"] if f["slide"] == 3]
    assert len(s3) == 1
    assert s3[0]["severity"] == "WARNING"


def test_exit_codes_gate(mixed_deck, clean_deck, capsys):
    assert audit_pptx.main([str(mixed_deck)]) == 1
    assert audit_pptx.main([str(clean_deck)]) == 0
    capsys.readouterr()


def test_strict_elevates_warning(clean_deck, tmp_path, capsys):
    prs = Presentation()
    _add_cjk_run(prs.slides.add_slide(prs.slide_layouts[6]), "只有继承警告")
    p = tmp_path / "warn.pptx"
    prs.save(p)
    assert audit_pptx.main([str(p)]) == 0
    assert audit_pptx.main([str(p), "--strict"]) == 1
    capsys.readouterr()


def test_sections_all_smoke(mixed_deck):
    report = audit_pptx.audit(mixed_deck, list(audit_pptx.SECTIONS))
    assert {*audit_pptx.SECTIONS} <= set(report)

    targets = [h["target"] for h in report["hyperlinks"]]
    assert "https://example.com/data" in targets

    assert report["metadata"]["core"].get("title") == "审计夹具"
    assert len(report["shapes"]) >= 3
    assert all("slide" in s and "kind" in s for s in report["shapes"])
    assert report["masters"]["masters"] == 1
    assert sum(report["masters"]["layout_usage"].values()) == 3
    assert report["security"] == {"msip_labels": [], "custom_xml_parts": []}
    assert report["themes"] and "colors" in report["themes"][0]


def test_bad_inputs(tmp_path, capsys):
    assert audit_pptx.main([str(tmp_path / "nope.pptx")]) == 2
    fake = tmp_path / "fake.pptx"
    fake.write_text("not a zip")
    assert audit_pptx.main([str(fake)]) == 2
    assert audit_pptx.main([str(fake), "--sections", "bogus"]) == 2
    capsys.readouterr()
