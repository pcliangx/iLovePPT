"""验证 build 主路径接通三层 SSOT(P0 修复):

1. load_theme(内置名)走 themes/_base.load_and_apply —— yaml token(色板/字体/
   mode/style)真实推到 module,ThemeConfig 挂 module._THEME_CONFIG
2. tier2 按 yaml layouts mapping 分发(alias:quadrant → make_matrix_2x2)
3. theme yaml 值为 null / 缺 key 的 layout fall back 到 helpers LayoutRegistry
   plugin 标准实现(17 enum 全部可渲染)
4. 三层都没有仍 fail-loud(tier3 不 silent remap)
"""
import json
from pathlib import Path

import pytest
from pptx import Presentation

from build import build_deck, load_plan, load_theme
from builder.tier2 import resolve_layout_fn
from themes import _base as theme_base


def _write_plan(tmp_path: Path, plan: dict) -> Path:
    p = tmp_path / "deck_v1_plan.json"
    p.write_text(json.dumps(plan, ensure_ascii=False), encoding="utf-8")
    return p


# ===========================================================================
# 1. yaml token 接线
# ===========================================================================

def test_load_theme_attaches_theme_config():
    mod = load_theme("tech_blue")
    cfg = getattr(mod, "_THEME_CONFIG", None)
    assert cfg is not None
    assert cfg.name == "tech_blue"


def test_load_theme_applies_yaml_tokens():
    """mode / style 配方常量必须被 apply_theme 推到 module(v0.15.0 特性,
    此前 build 路径不可达)。"""
    mod = load_theme("tech_blue")
    assert mod.THEME_MODE == "light"
    assert mod.STYLE_RECIPE in theme_base.STYLE_RECIPES
    assert isinstance(mod.STYLE_RADIUS_MEDIUM, float)
    assert isinstance(mod.STYLE_GAP_IN, float)


def test_load_theme_yaml_color_matches_module_primary():
    """yaml colors.brand_primary 与 module.PRIMARY 一致(token 单源)。"""
    mod = load_theme("tech_blue")
    cfg = mod._THEME_CONFIG
    assert mod.PRIMARY == cfg.colors["brand_primary"]
    assert mod.ACCENT == cfg.colors["accent"]


def test_load_theme_template_golden_now_loadable():
    """template_golden 有 yaml,应可直接作为 build theme(此前只走 .pptx 提取旁路)。"""
    mod = load_theme("template_golden")
    assert mod._THEME_CONFIG.name == "template_golden"


# ===========================================================================
# 2. yaml alias 分发
# ===========================================================================

def test_quadrant_alias_dispatches_to_matrix_2x2():
    mod = load_theme("tech_blue")
    fn, is_plugin = resolve_layout_fn(mod, "quadrant")
    assert fn is mod.make_matrix_2x2
    assert is_plugin is False


def test_comparison_alias_dispatches_to_compare():
    mod = load_theme("tech_blue")
    fn, is_plugin = resolve_layout_fn(mod, "comparison")
    assert fn is mod.make_compare
    assert is_plugin is False


def test_build_deck_renders_quadrant_via_alias(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [{
            "layout": "quadrant", "title": "定位矩阵",
            "x_axis": {"low": "低成本", "high": "高成本"},
            "y_axis": {"low": "低价值", "high": "高价值"},
            "quadrants": [
                {"pos": "tl", "title": "Q1", "desc": "d1"},
                {"pos": "tr", "title": "Q2", "desc": "d2"},
                {"pos": "bl", "title": "Q3", "desc": "d3"},
                {"pos": "br", "title": "Q4", "desc": "d4"},
            ]}]})
    out = build_deck(load_plan(p))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


# ===========================================================================
# 3. plugin fallback(yaml null / 缺 key → LayoutRegistry)
# ===========================================================================

def test_null_layout_falls_back_to_plugin():
    """tech_blue.yaml layouts.data: null → helpers/data.py plugin 兜底。"""
    mod = load_theme("tech_blue")
    fn, is_plugin = resolve_layout_fn(mod, "data")
    assert fn is not None
    assert is_plugin is True


@pytest.mark.parametrize("layout", [
    "data", "timeline", "pyramid", "venn", "radial", "process_flow", "quote",
])
def test_all_null_layouts_resolvable_under_tech_blue(layout):
    """17 enum 里 tech_blue 未实现的 layout 全部有 plugin 兜底 —— author 按
    受控词典选任何 layout_type 都不会 build 期撞墙。"""
    mod = load_theme("tech_blue")
    fn, is_plugin = resolve_layout_fn(mod, layout)
    assert fn is not None, f"layout {layout!r} 三层都解析不到"
    assert is_plugin is True


def test_build_deck_renders_timeline_via_plugin(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [{
            "layout": "timeline", "title": "里程碑",
            "milestones": [
                {"title": "启动", "desc": "kickoff", "date": "Q1"},
                {"title": "交付", "desc": "ship", "date": "Q3"},
            ]}]})
    out = build_deck(load_plan(p))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_build_deck_renders_data_via_plugin(tmp_path):
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [{
            "layout": "data", "title": "季度数据",
            "headers": ["指标", "Q1", "Q2"],
            "table_rows": [["营收", "1.2", "1.5"], ["毛利", "0.4", "0.6"]]}]})
    out = build_deck(load_plan(p))
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_plugin_page_gets_footer(tmp_path):
    """plugin 渲染的内容页也要有页脚页码(FOOTERED_LAYOUTS 已补齐)。"""
    p = _write_plan(tmp_path, {
        "theme": "tech_blue", "output": "./d.pptx",
        "slides": [
            {"layout": "cover", "title": "T", "subtitle": "S"},
            {"layout": "data", "title": "数据",
             "headers": ["A"], "table_rows": [["1"]]}]})
    out = build_deck(load_plan(p))
    prs = Presentation(str(out))
    texts = [sh.text_frame.text for sh in prs.slides[1].shapes
             if sh.has_text_frame]
    assert any("1 / 1" in t for t in texts), f"plugin 页缺页码,texts={texts}"


# ===========================================================================
# 4. fail-loud 保持
# ===========================================================================

def test_unknown_layout_still_fails_loud():
    mod = load_theme("tech_blue")
    fn, is_plugin = resolve_layout_fn(mod, "nonexistent_xyz")
    assert fn is None


def test_yaml_declared_but_missing_impl_raises():
    """yaml 声明了函数名但 module 没有 → ValueError(不静默降级 plugin)。"""
    mod = load_theme("tech_blue")
    cfg = mod._THEME_CONFIG
    original = cfg.layouts.get("cover")
    cfg.layouts["cover"] = "make_does_not_exist"
    try:
        with pytest.raises(ValueError, match="make_does_not_exist"):
            resolve_layout_fn(mod, "cover")
    finally:
        cfg.layouts["cover"] = original
