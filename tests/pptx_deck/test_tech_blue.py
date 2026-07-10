# tests/pptx_deck/test_tech_blue.py
"""tech_blue 主题 13 layout light test：验证每个 layout 创建后 prs.slides 增加 1。"""
from pptx import Presentation
from pptx.util import Inches
import helpers as H
from themes import tech_blue as T


def _new():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def test_make_cover():
    prs = _new()
    T.make_cover(prs, "主标题", "副标题")
    assert len(prs.slides) == 1

def test_make_toc():
    prs = _new()
    T.make_toc(prs, sections=["背景", "范围", "流程", "保障", "节奏"])
    assert len(prs.slides) == 1

def test_make_section_divider():
    prs = _new()
    T.make_section_divider(prs, 1, "第一章")
    assert len(prs.slides) == 1

def test_make_single_focus():
    prs = _new()
    T.make_single_focus(prs, big_text="一句话", big_number="80%", explanation="解释")
    assert len(prs.slides) == 1

def test_make_compare_two_items():
    prs = _new()
    T.make_compare(prs, "对比", items=[
        {"title": "现状", "body": "手工"},
        {"title": "目标", "body": "自动"}])
    assert len(prs.slides) == 1

def test_make_compare_three_items():
    prs = _new()
    T.make_compare(prs, "三方对比", items=[
        {"title": "A", "body": "x"}, {"title": "B", "body": "y"},
        {"title": "C", "body": "z"}])
    assert len(prs.slides) == 1

def test_make_compare_with_recommended_adds_badge():
    """recommended=True 主推列 = 多 1 个 OVAL 徽章 + 1 个徽章文字 = 2 shapes vs 普通列。"""
    prs = _new()
    T.make_compare(prs, "高亮", items=[
        {"title": "普通", "body": "x"},
        {"title": "主推", "body": "y", "recommended": True}])
    assert len(prs.slides) == 1
    # 每列基础 4 shapes(header rect + header textbox + body rect + body textbox)
    # 主推列多 1 OVAL + 1 textbox(徽章)→ title 1 + 2*4 + 2 = 11
    assert len(prs.slides[0].shapes) == 11

def test_make_compare_pk():
    prs = _new()
    T.make_compare_pk(prs, "对决",
                      left={"title": "旧", "body": "旧方案 body"},
                      right={"title": "新", "body": "新方案 body"})
    assert len(prs.slides) == 1
    # title + 2 sides × (bg + bar + title + body) + VS circle + VS text = 1+8+2 = 11
    assert len(prs.slides[0].shapes) == 11

def test_make_matrix_2x2():
    prs = _new()
    T.make_matrix_2x2(prs, "矩阵",
        x_axis={"low": "x 低", "high": "x 高"},
        y_axis={"low": "y 低", "high": "y 高"},
        quadrants=[
            {"pos": "tl", "title": "tl", "body": "tl body"},
            {"pos": "tr", "title": "主推", "body": "tr body", "highlight": True},
            {"pos": "bl", "title": "bl", "body": "bl body"},
            {"pos": "br", "title": "br", "body": "br body"}])
    assert len(prs.slides) == 1
    # title + 4 象限 × (rect + title textbox + body textbox) + 4 axis labels = 1+12+4 = 17
    assert len(prs.slides[0].shapes) == 17

def test_make_matrix_2x2_invalid_pos_raises():
    prs = _new()
    import pytest
    with pytest.raises(ValueError, match="quadrant.pos"):
        T.make_matrix_2x2(prs, "x",
            x_axis={"low": "a", "high": "b"},
            y_axis={"low": "c", "high": "d"},
            quadrants=[{"pos": "xx", "title": "t", "body": "b"}])

def test_make_cards_two():
    prs = _new()
    T.make_cards(prs, "两栏", cards=[
        {"title": "卡1", "body": "正文1"}, {"title": "卡2", "body": "正文2"}])
    assert len(prs.slides) == 1

def test_make_cards_four():
    prs = _new()
    T.make_cards(prs, "四栏", cards=[
        {"title": f"卡{i}", "body": f"正文{i}"} for i in range(1, 5)])
    assert len(prs.slides) == 1

def test_make_bullet_list():
    prs = _new()
    T.make_bullet_list(prs, "标题", items=["要点1", "要点2", "要点3", "要点4", "要点5"])
    assert len(prs.slides) == 1
    # 短条目走"逐条整页分布"路径:title + 每条(accent rect + 文字框)。
    # 长条目塞满时退回紧凑 H.bullets(title + 单 bullets 框)。两种路径都 >= 2 shapes。
    shapes = prs.slides[0].shapes
    assert len(shapes) >= 2  # 至少 title + 内容


def test_make_bullet_list_long_items_no_overflow():
    """长句条目(handout)塞满时退回紧凑路径,内容不溢出页脚(y < FOOTER_TOP)。"""
    prs = _new()
    H.PRESENTATION_MODE = "handout"
    try:
        long_items = [
            "城市等级与学历卡方 chi2=704、p=2.8e-140,极度显著,下沉买家学历确实更低,城市与学历强耦合、是真实结构。",
            "城市等级与年龄卡方 chi2=44、p=0.044,仅弱显著,远不及城市与学历的耦合强度,四象限的年龄轴解释力有限。",
            "两组卡方印证:城市、学历这类硬人口学维度是可信真信号,与前页均匀偏置的推断标签形成鲜明对比。",
            "边界声明:聚类 silhouette 仅 0.08–0.10,买家同质、无泾渭分明的自然客群,打法差异宜基于象限、不宜过度精细化。",
        ]
        T.make_bullet_list(prs, "卡方证实", items=long_items)
        # 所有文字框底边不得越过 FOOTER_TOP(7.0in)— source 引文 y=6.7in 之上更佳
        for sh in prs.slides[0].shapes:
            if sh.has_text_frame and sh.top is not None and sh.height is not None:
                assert sh.top + sh.height <= H.FOOTER_TOP + 1, (
                    f"shape bottom {(sh.top+sh.height)/914400:.2f}in 越过 footer"
                )
    finally:
        H.PRESENTATION_MODE = "speaker"

def test_make_table():
    prs = _new()
    T.make_table(prs, "表格标题",
                 headers=["A", "B", "C"],
                 rows=[["1", "2", "3"], ["4", "5", "6"]])
    assert len(prs.slides) == 1
    assert len(prs.slides[0].shapes) == 2  # title textbox + table

def test_make_pic_text(tmp_path):
    from PIL import Image
    img_path = tmp_path / "blank.png"
    Image.new("RGB", (10, 10), "white").save(str(img_path))
    prs = _new()
    T.make_pic_text(prs, "标题", str(img_path),
                    points=[{"title": "点1", "body": "正文1"},
                            {"title": "点2", "body": "正文2"}])
    assert len(prs.slides) == 1

def test_make_summary():
    prs = _new()
    T.make_summary(prs, conclusions=["结论 1", "结论 2", "结论 3"])
    assert len(prs.slides) == 1

def test_make_closing():
    prs = _new()
    T.make_closing(prs, subtitle="联系邮箱：x@y.com")
    assert len(prs.slides) == 1

def test_font_default_is_microsoft_yahei():
    assert T.FONT_HEADER == "Microsoft YaHei"
    assert T.FONT_BODY == "Microsoft YaHei"


# ----- 2026-05-23 新增字段测试 -----

def test_make_cover_with_metadata():
    """cover 接受 prepared_by / date / version / project_code / classification。"""
    prs = _new()
    T.make_cover(prs, "T", "S",
                 prepared_by="技术部",
                 date="2026-05-23",
                 version="v1.0",
                 project_code="ATLAS-01",
                 classification="INTERNAL")
    slide = prs.slides[0]
    texts = [sh.text_frame.text for sh in slide.shapes
             if sh.has_text_frame and sh.text_frame.text]
    joined = " | ".join(texts)
    assert "INTERNAL" in joined
    assert "技术部" in joined
    assert "2026-05-23" in joined
    assert "v1.0" in joined
    assert "ATLAS-01" in joined


def test_make_cover_no_metadata_still_works():
    """cover 不传任何 meta 字段时,保持原有 title + subtitle 渲染。"""
    prs = _new()
    T.make_cover(prs, "T", "S")
    slide = prs.slides[0]
    texts = [sh.text_frame.text for sh in slide.shapes
             if sh.has_text_frame and sh.text_frame.text]
    joined = " | ".join(texts)
    assert "T" in joined and "S" in joined


def test_make_closing_with_next_steps():
    """next_steps 模式:渲染 'Next Steps' 标题 + 编号 action 列表。"""
    prs = _new()
    T.make_closing(prs, subtitle="问答 / 联系",
                   next_steps=[
                       {"action": "完成 Phase 1 试点", "owner": "Alice",
                        "due": "2026-06-15"},
                       {"action": "评估扩展到 Phase 2", "owner": "Bob",
                        "due": "2026-07-01"}])
    slide = prs.slides[0]
    texts = " | ".join(sh.text_frame.text for sh in slide.shapes
                       if sh.has_text_frame and sh.text_frame.text)
    assert "Next Steps" in texts
    assert "完成 Phase 1 试点" in texts
    assert "Alice" in texts and "2026-06-15" in texts
    assert "评估扩展到 Phase 2" in texts


def test_make_closing_without_next_steps_uses_simple_mode():
    """无 next_steps 时退回原版'谢谢'封底。"""
    prs = _new()
    T.make_closing(prs, subtitle="github.com/xxx")
    slide = prs.slides[0]
    texts = " | ".join(sh.text_frame.text for sh in slide.shapes
                       if sh.has_text_frame and sh.text_frame.text)
    assert "谢谢" in texts
    assert "Next Steps" not in texts


def test_make_closing_headline_mode():
    """headline 模式:大字主句居中 + note 小字副行,无硬编码 'Next Steps' / '谢谢'。"""
    prs = _new()
    T.make_closing(prs, headline="数据揭示问题,对策由一线研判",
                   note="可复现源 — analysis/ 脚本 + processed CSV")
    slide = prs.slides[0]
    texts = " | ".join(sh.text_frame.text for sh in slide.shapes
                       if sh.has_text_frame and sh.text_frame.text)
    assert "数据揭示问题,对策由一线研判" in texts  # headline 主句
    assert "可复现源" in texts                       # note 副行
    assert "Next Steps" not in texts                 # 不再硬编码 Next Steps
    assert "谢谢" not in texts                        # headline 优先于简单模式
    # headline 文字框应垂直居中(vertical_anchor=MIDDLE=3),避免下沉/空白
    head_boxes = [sh for sh in slide.shapes
                  if sh.has_text_frame and "数据揭示问题" in sh.text_frame.text]
    assert head_boxes
    from pptx.enum.text import MSO_ANCHOR
    assert head_boxes[0].text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE


def test_make_closing_headline_overrides_next_steps():
    """headline 优先级最高:即使误传 next_steps,也走极简大字、不渲 Next Steps 列表。"""
    prs = _new()
    T.make_closing(prs, headline="一句话主张",
                   next_steps=[{"action": "不应出现"}])
    slide = prs.slides[0]
    texts = " | ".join(sh.text_frame.text for sh in slide.shapes
                       if sh.has_text_frame and sh.text_frame.text)
    assert "一句话主张" in texts
    assert "Next Steps" not in texts
    assert "不应出现" not in texts
